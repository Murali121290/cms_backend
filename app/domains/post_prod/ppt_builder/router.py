import os
import shutil
import json
import re
import time
import fitz  # PyMuPDF
from pathlib import Path
from typing import Optional
from fastapi import APIRouter, UploadFile, File, Form, HTTPException, Cookie, Header, Depends
from fastapi.responses import FileResponse, JSONResponse
from pptx import Presentation
from pptx.util import Pt

from app.domains.auth.security import get_current_user_from_cookie
from app.domains.auth.rbac_config import has_post_prod_access

from .services.main import extract_template
from .services.convert import convert
from .services.excel_report import create_excel_report
from .services.report import collect_changes, collect_figure_diagnostics
from .services.accessibility import check_ppt_accessibility

def check_post_prod_access(user = Depends(get_current_user_from_cookie)):
    if not user or not has_post_prod_access(user):
        raise HTTPException(status_code=403, detail="Access denied to Post Production / Backlist.")
    return user

router = APIRouter(
    prefix="/post-prod/ppt-builder",
    tags=["PPT Builder"],
    dependencies=[Depends(check_post_prod_access)]
)

UPLOAD_DIR = os.path.abspath("uploads/ppt_builder")
TEMPLATES_DIR = os.path.join(UPLOAD_DIR, "templates")
os.makedirs(TEMPLATES_DIR, exist_ok=True)

# In-memory store for session-specific state
session_states = {}

def get_session_id(session_id: str = Cookie(default="default"), x_session_id: str = Header(default=None)):
    if x_session_id:
        return x_session_id
    return session_id

def get_session_state(session_id: str) -> dict:
    if session_id not in session_states:
        session_states[session_id] = {
            "template_pptx": None,
            "template_style_json": None,
            "content_pptx": None,
            "styled_pptx": None,
            "pdf_path": None,
            "captions": []
        }
    return session_states[session_id]

def reset_session_state(session_id: str):
    session_states[session_id] = {
        "template_pptx": None,
        "template_style_json": None,
        "content_pptx": None,
        "styled_pptx": None,
        "pdf_path": None,
        "captions": []
    }

def cleanup_old_sessions():
    now = time.time()
    cutoff = now - 24 * 3600  # 24 hours ago
    if not os.path.exists(UPLOAD_DIR):
        return
    for item in os.listdir(UPLOAD_DIR):
        item_path = os.path.join(UPLOAD_DIR, item)
        if item == "templates":
            continue
        if os.path.isdir(item_path):
            try:
                mtime = os.path.getmtime(item_path)
                if mtime < cutoff:
                    shutil.rmtree(item_path, ignore_errors=True)
                    if item in session_states:
                        del session_states[item]
            except Exception as e:
                print(f"Failed to check/cleanup old session folder {item}: {e}")

def get_session_upload_dir(session_id: str) -> str:
    try:
        cleanup_old_sessions()
    except Exception:
        pass
    clean_id = re.sub(r'[^a-zA-Z0-9_-]', '', session_id)
    if not clean_id:
        clean_id = "default"
    path = os.path.join(UPLOAD_DIR, clean_id)
    os.makedirs(path, exist_ok=True)
    os.makedirs(os.path.join(path, "pdf_extracts"), exist_ok=True)
    os.makedirs(os.path.join(path, "rendered_slides"), exist_ok=True)
    return path

def _span_style(span):
    flags = span.get("flags", 0)
    font_name = span.get("font", "").lower()
    bold = bool(flags & 16) or "bold" in font_name
    italic = bool(flags & 2) or "italic" in font_name or "oblique" in font_name
    return bold, italic

def _block_runs(block):
    runs = []
    for line in block.get("lines", []):
        for span in line.get("spans", []):
            text = span.get("text", "")
            if text:
                bold, italic = _span_style(span)
                runs.append({"text": text, "bold": bold, "italic": italic})
    return runs

def _block_plain(block):
    return " ".join(
        span.get("text", "")
        for line in block.get("lines", [])
        for span in line.get("spans", [])
    ).strip()

def extract_pdf_captions(pdf_path):
    captions = []
    pattern = re.compile(r"^\s*(Figure|Fig\.|Table)\s+(\d+[-.\d]*)\b", re.IGNORECASE)
    credit_pattern = re.compile(
        r"^\s*(©|Copyright\b|Courtesy of\b|Source:\b|Source\b|Reproduced from\b|"
        r"Reproduced with permission\b|Data from\b|Courtesy\b|Permission\b|Used with permission\b)",
        re.IGNORECASE,
    )

    if not os.path.exists(pdf_path):
        return captions

    try:
        doc = fitz.open(pdf_path)
        for page_idx in range(len(doc)):
            page = doc[page_idx]
            page_data = page.get_text("dict")
            blocks = [b for b in page_data.get("blocks", []) if b.get("type") == 0]

            for b_idx, block in enumerate(blocks):
                plain = _block_plain(block)
                if not plain:
                    continue

                lines = []
                for line in block.get("lines", []):
                    line_text = " ".join(s.get("text", "") for s in line.get("spans", [])).strip()
                    if line_text:
                        lines.append(line_text)

                for line_idx, line in enumerate(lines):
                    match = pattern.match(line.strip())
                    if match:
                        full_caption = " ".join(lines[line_idx:])
                        clean_text = " ".join(full_caption.split())
                        cap_runs = _block_runs(block)

                        credit_text = ""
                        credit_runs = []
                        for nb in blocks[b_idx + 1: b_idx + 3]:
                            nb_plain = _block_plain(nb)
                            if credit_pattern.match(nb_plain):
                                credit_text = " ".join(nb_plain.split())
                                credit_runs = _block_runs(nb)
                            break

                        cap_type = match.group(1).capitalize()
                        if cap_type.startswith("Fig"):
                            cap_type = "Figure"

                        captions.append({
                            "id": f"cap_{page_idx}_{b_idx}_{line_idx}",
                            "page": page_idx + 1,
                            "label": f"{cap_type} {match.group(2)}",
                            "text": clean_text,
                            "runs": cap_runs,
                            "credit": credit_text,
                            "creditRuns": credit_runs,
                        })
                        break
        doc.close()
    except Exception as e:
        print("Error extracting PDF captions:", e)
    return captions

@router.post("/reset")
async def reset_session(session_id: str = Depends(get_session_id)):
    try:
        session_upload_dir = get_session_upload_dir(session_id)
        for item in os.listdir(session_upload_dir):
            item_path = os.path.join(session_upload_dir, item)
            if item == "templates":
                continue
            if os.path.isdir(item_path):
                shutil.rmtree(item_path, ignore_errors=True)
            else:
                try:
                    os.remove(item_path)
                except Exception:
                    pass
        os.makedirs(os.path.join(session_upload_dir, "pdf_extracts"), exist_ok=True)
        os.makedirs(os.path.join(session_upload_dir, "rendered_slides"), exist_ok=True)
        reset_session_state(session_id)
        return {"ok": True, "message": "Session reset. All uploaded files cleared."}
    except Exception as e:
        return {"ok": False, "detail": str(e)}

@router.post("/upload-template")
async def upload_template(file: UploadFile = File(...), session_id: str = Depends(get_session_id)):
    try:
        filename = file.filename
        path = os.path.join(TEMPLATES_DIR, filename)
        with open(path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        state = get_session_state(session_id)
        state["template_pptx"] = path
        
        styles = extract_template(path)
        style_json_filename = os.path.splitext(filename)[0] + "_styles.json"
        style_json_path = os.path.join(TEMPLATES_DIR, style_json_filename)
        with open(style_json_path, "w") as f:
            json.dump(styles, f, indent=2)
            
        state["template_style_json"] = style_json_path
        return {"ok": True, "styles": styles, "filename": filename}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/templates")
async def get_templates(session_id: str = Depends(get_session_id)):
    try:
        templates = []
        if os.path.exists(TEMPLATES_DIR):
            for f in os.listdir(TEMPLATES_DIR):
                if f.endswith(".pptx"):
                    templates.append({
                        "name": os.path.splitext(f)[0],
                        "filename": f
                    })
        return {"ok": True, "templates": templates}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/select-template")
async def select_template(data: dict, session_id: str = Depends(get_session_id)):
    filename = data.get("filename")
    if not filename:
        raise HTTPException(status_code=400, detail="Missing filename")
    
    path = os.path.join(TEMPLATES_DIR, filename)
    style_json_filename = os.path.splitext(filename)[0] + "_styles.json"
    style_json_path = os.path.join(TEMPLATES_DIR, style_json_filename)
    
    if not os.path.exists(path) or not os.path.exists(style_json_path):
        raise HTTPException(status_code=404, detail="Template or styles not found")
        
    state = get_session_state(session_id)
    state["template_pptx"] = path
    state["template_style_json"] = style_json_path
    
    with open(style_json_path, "r") as f:
        styles = json.load(f)
        
    return {"ok": True, "filename": filename, "styles": styles}

@router.post("/upload-ppt")
async def upload_ppt(file: UploadFile = File(...), session_id: str = Depends(get_session_id)):
    try:
        session_upload_dir = get_session_upload_dir(session_id)
        path = os.path.join(session_upload_dir, "uploaded_content.pptx")
        with open(path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        state = get_session_state(session_id)
        state["content_pptx"] = path
        
        slides_info = extract_template(path)
        
        return {"ok": True, "filename": file.filename, "slidesInfo": slides_info}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/process-ppt")
async def process_ppt(payload: dict = None, session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    if not state["content_pptx"] or not state["template_style_json"]:
        raise HTTPException(status_code=400, detail="Missing uploaded PPTX or Template style JSON")
    
    try:
        session_upload_dir = get_session_upload_dir(session_id)
        extracts_dir = os.path.join(session_upload_dir, "pdf_extracts")
        if os.path.exists(extracts_dir):
            shutil.rmtree(extracts_dir)
        os.makedirs(extracts_dir, exist_ok=True)

        include_figure_captions = (payload or {}).get("include_figure_captions", True)
        include_table_captions = (payload or {}).get("include_table_captions", True)

        figures = (payload or {}).get("figures", [])
        for fig in figures:
            name = fig.get("name")
            filename = fig.get("filename")
            if name and filename:
                src_path = os.path.join(session_upload_dir, filename)
                if os.path.exists(src_path):
                    m = re.match(r"(figure|table)\s*([\d.-]+)", name, re.IGNORECASE)
                    if m:
                        dest_name = f"{m.group(1).lower()} {m.group(2)}.png"
                    else:
                        dest_name = f"{name.lower()}.png"
                    shutil.copy(src_path, os.path.join(extracts_dir, dest_name))

        output_path = os.path.join(session_upload_dir, "styled_output.pptx")
        used_figs = convert(
            state["content_pptx"], state["template_style_json"], output_path,
            apply_geometry=True, cover_image_path=None, figures_metadata=figures,
            include_figure_captions=include_figure_captions, include_table_captions=include_table_captions
        )
        state["styled_pptx"] = output_path

        auto_inserted_list = []
        for uf in used_figs:
            dest_name = uf["dest_name"]
            orig_filename = None
            for fig in figures:
                name = fig.get("name")
                filename = fig.get("filename")
                if name and filename:
                    m = re.match(r"(figure|table)\s*([\d.-]+)", name, re.IGNORECASE)
                    if m:
                        candidate = f"{m.group(1).lower()} {m.group(2)}.png"
                    else:
                        candidate = f"{name.lower()}.png"
                    if candidate == dest_name:
                        orig_filename = filename
                        break
            if orig_filename:
                auto_inserted_list.append({
                    "filename": orig_filename,
                    "slideIndex": uf["slideIndex"],
                    "shapeIndex": uf["shapeIndex"]
                })
        
        state["auto_inserted_list"] = auto_inserted_list
        slides_info = extract_template(output_path)
        
        prs = Presentation(output_path)
        extracted_images = {}
        for s_idx, slide in enumerate(prs.slides):
            for sh_idx, shape in enumerate(slide.shapes):
                if hasattr(shape, "shape_type") and int(shape.shape_type) == 13:
                    try:
                        img = shape.image
                        ext = img.ext
                        img_filename = f"slide_{s_idx}_shape_{sh_idx}.{ext}"
                        img_path = os.path.join(session_upload_dir, img_filename)
                        with open(img_path, "wb") as f:
                            f.write(img.blob)
                        extracted_images[f"{s_idx}_{sh_idx}"] = f"/api/v2/post-prod/ppt-builder/media/{session_id}/{img_filename}"
                    except Exception:
                        pass
        
        for s_idx, slide_data in enumerate(slides_info.get("slides", [])):
            for sh_idx, shape_data in enumerate(slide_data.get("shapes", [])):
                key = f"{s_idx}_{sh_idx}"
                if key in extracted_images:
                    shape_data["imageUrl"] = extracted_images[key]

        return {
            "ok": True, 
            "slidesInfo": slides_info, 
            "autoInserted": state.get("auto_inserted_list", [])
        }
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/upload-pdf")
async def upload_pdf(file: UploadFile = File(...), session_id: str = Depends(get_session_id)):
    try:
        session_upload_dir = get_session_upload_dir(session_id)
        path = os.path.join(session_upload_dir, "uploaded_document.pdf")
        with open(path, "wb") as buffer:
            shutil.copyfileobj(file.file, buffer)
        
        state = get_session_state(session_id)
        state["pdf_path"] = path
        
        doc = fitz.open(path)
        page_count = len(doc)
        filename = file.filename
        doc.close()
        
        captions = extract_pdf_captions(path)
        state["captions"] = captions
        
        return {"ok": True, "filename": filename, "pageCount": page_count, "captions": captions}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/parse-alttext-excel")
async def parse_alttext_excel(file: UploadFile = File(...)):
    try:
        import openpyxl
        wb = openpyxl.load_workbook(file.file, data_only=True)
        sheet = wb.active
        
        entries = []
        headers = {}
        # Scan first row for columns
        for col_idx in range(1, 15):
            val = sheet.cell(row=1, column=col_idx).value
            if val:
                headers[str(val).strip().lower()] = col_idx
                
        def get_col_val(row_cells, names, default_idx):
            for name in names:
                if name in headers:
                    return row_cells[headers[name] - 1].value
            if len(row_cells) >= default_idx:
                return row_cells[default_idx - 1].value
            return None

        # Process data rows
        for r_idx in range(2, sheet.max_row + 1):
            row_cells = [sheet.cell(row=r_idx, column=c) for c in range(1, 10)]
            if not any(cell.value for cell in row_cells):
                continue
            
            fig_key = get_col_val(row_cells, ["figure key", "figure_key", "figure number", "figure no", "figure", "id"], 1)
            element = get_col_val(row_cells, ["element", "element number", "element no"], 2)
            chapter = get_col_val(row_cells, ["chapter"], 3)
            deco_val = get_col_val(row_cells, ["decorative", "is_decorative"], 4)
            alt_short = get_col_val(row_cells, ["alt text short", "alt_text_short", "short alt text", "alt text", "alttext", "description"], 5)
            alt_long = get_col_val(row_cells, ["alt text long", "alt_text_long", "long alt text", "caption"], 6)
            
            if not fig_key:
                continue
                
            entries.append({
                "figure_key": str(fig_key).strip().lower(),
                "element": str(element or ""),
                "chapter": str(chapter or ""),
                "decorative": bool(deco_val),
                "alt_text_short": str(alt_short or ""),
                "alt_text_long": str(alt_long or "")
            })
            
        return {"ok": True, "entries": entries}
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to parse alt-text excel: {str(e)}")

@router.get("/pdf/captions")
async def get_pdf_captions(session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    if not state["pdf_path"]:
        raise HTTPException(status_code=404, detail="No PDF uploaded")
    if not state.get("captions"):
        state["captions"] = extract_pdf_captions(state["pdf_path"])
    return {"ok": True, "captions": state["captions"]}

@router.get("/pdf/info")
async def get_pdf_info(session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    if not state["pdf_path"]:
        raise HTTPException(status_code=404, detail="No PDF uploaded")
    try:
        doc = fitz.open(state["pdf_path"])
        count = len(doc)
        filename = os.path.basename(state["pdf_path"])
        doc.close()
        return {"filename": filename, "count": count}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/pdf/file")
async def get_pdf_file(session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    if not state["pdf_path"]:
        raise HTTPException(status_code=404, detail="No PDF uploaded")
    return FileResponse(state["pdf_path"], media_type="application/pdf")

@router.post("/extract")
async def extract_crop(
    page: int = Form(...),
    scale: float = Form(...),
    x0: float = Form(...),
    y0: float = Form(...),
    x1: float = Form(...),
    y1: float = Form(...),
    session_id: str = Depends(get_session_id)
):
    state = get_session_state(session_id)
    if not state["pdf_path"]:
        raise HTTPException(status_code=404, detail="No PDF uploaded")
    try:
        session_upload_dir = get_session_upload_dir(session_id)
        doc = fitz.open(state["pdf_path"])
        pdf_page = doc[page]
        
        clip = fitz.Rect(
            min(x0, x1) / scale,
            min(y0, y1) / scale,
            max(x0, x1) / scale,
            max(y0, y1) / scale
        )
        
        pix = pdf_page.get_pixmap(
            matrix=fitz.Matrix(3.0, 3.0),
            clip=clip,
            alpha=False
        )
        png_data = pix.tobytes("png")
        doc.close()
        
        existing = [f for f in os.listdir(session_upload_dir) if f.startswith("crop_") and f.endswith(".png")]
        idx = len(existing) + 1
        filename = f"crop_p{page+1}_{idx:03d}.png"
        filepath = os.path.join(session_upload_dir, filename)
        
        with open(filepath, "wb") as f:
            f.write(png_data)
            
        return {"filename": filename, "url": f"/api/v2/post-prod/ppt-builder/media/{session_id}/{filename}", "page": page + 1}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.post("/add-image")
async def add_image(
    slide_index: int = Form(...),
    image_name: str = Form(...),
    x_pt: float = Form(...),
    y_pt: float = Form(...),
    w_pt: float = Form(...),
    h_pt: float = Form(...),
    caption: str = Form(None),
    caption_runs: str = Form(None),
    session_id: str = Depends(get_session_id)
):
    state = get_session_state(session_id)
    target_pptx = state["styled_pptx"] or state["content_pptx"]
    if not target_pptx:
        raise HTTPException(status_code=400, detail="No PPTX presentation available to modify")
        
    session_upload_dir = get_session_upload_dir(session_id)
    extracts_dir = os.path.join(session_upload_dir, "pdf_extracts")
    image_path = os.path.join(extracts_dir, image_name)
    if not os.path.exists(image_path):
        image_path = os.path.join(session_upload_dir, image_name)
    if not os.path.exists(image_path):
        raise HTTPException(status_code=404, detail=f"Image {image_name} not found")
        
    try:
        prs = Presentation(target_pptx)
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise HTTPException(status_code=400, detail="Invalid slide index")
            
        slide = prs.slides[slide_index]
        
        left = Pt(x_pt)
        top = Pt(y_pt)
        width = Pt(w_pt)
        height = Pt(h_pt)
        
        pic = slide.shapes.add_picture(image_path, left, top, width, height)
        
        if caption:
            cap_top = top + height + Pt(10)
            cap_height = Pt(35)
            if cap_top + cap_height > prs.slide_height:
                cap_top = prs.slide_height - cap_height - Pt(10)
            txBox = slide.shapes.add_textbox(left, cap_top, width, cap_height)
            tf = txBox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            runs_data = None
            if caption_runs:
                try:
                    runs_data = json.loads(caption_runs)
                except Exception:
                    pass
            if runs_data:
                for run_data in runs_data:
                    run = p.add_run()
                    run.text = run_data.get("text", "")
                    run.font.size = Pt(10)
                    run.font.bold = run_data.get("bold", False)
                    run.font.italic = run_data.get("italic", True)
            else:
                run = p.add_run()
                run.text = caption
                run.font.size = Pt(10)
                run.font.italic = True
            
        prs.save(target_pptx)
        slides_info = extract_template(target_pptx)
        
        extracted_images = {}
        for s_idx, s in enumerate(prs.slides):
            for sh_idx, shape in enumerate(s.shapes):
                if hasattr(shape, "shape_type") and int(shape.shape_type) == 13:
                    try:
                        img = shape.image
                        ext = img.ext
                        img_filename = f"slide_{s_idx}_shape_{sh_idx}.{ext}"
                        img_path = os.path.join(session_upload_dir, img_filename)
                        with open(img_path, "wb") as f:
                            f.write(img.blob)
                        extracted_images[f"{s_idx}_{sh_idx}"] = f"/api/v2/post-prod/ppt-builder/media/{session_id}/{img_filename}"
                    except Exception:
                        pass
                        
        for s_idx, slide_data in enumerate(slides_info.get("slides", [])):
            for sh_idx, shape_data in enumerate(slide_data.get("shapes", [])):
                key = f"{s_idx}_{sh_idx}"
                if key in extracted_images:
                    shape_data["imageUrl"] = extracted_images[key]
                    
        return {"ok": True, "slidesInfo": slides_info}
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))

@router.get("/download")
async def download_pptx(session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    target_pptx = state["styled_pptx"] or state["content_pptx"]
    if not target_pptx or not os.path.exists(target_pptx):
        raise HTTPException(status_code=404, detail="No output PPTX available to download")
    return FileResponse(
        target_pptx,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        filename="styled_presentation.pptx"
    )

@router.get("/download-excel")
async def download_excel(customerName: str = "", projectName: str = "", session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    input_path = state.get("content_pptx")
    session_upload_dir = get_session_upload_dir(session_id)
    output_path = state.get("styled_pptx") or os.path.join(session_upload_dir, "styled_output.pptx")
    
    if not input_path or not os.path.exists(input_path):
        raise HTTPException(status_code=400, detail="Missing source presentation file.")
    if not os.path.exists(output_path):
        output_path = input_path
        
    excel_path = os.path.join(session_upload_dir, "compilation_report.xlsx")
    extracts_dir = os.path.join(session_upload_dir, "pdf_extracts")
    
    try:
        create_excel_report(
            input_pptx=input_path,
            output_pptx=output_path,
            extracts_dir=extracts_dir,
            customer_name=customerName,
            project_name=projectName,
            output_excel_path=excel_path
        )
        
        return FileResponse(
            excel_path,
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename="compilation_report.xlsx"
        )
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Failed to generate Excel report: {str(e)}")

@router.get("/report-data")
async def get_report_data(session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    session_upload_dir = get_session_upload_dir(session_id)
    output_path = os.path.join(session_upload_dir, "styled_output.pptx")
    if not state["content_pptx"] or not os.path.exists(output_path):
        return {"ok": False, "changes": []}
    try:
        changes = collect_changes(state["content_pptx"], output_path)
        return {"ok": True, "changes": changes}
    except Exception as e:
        return {"ok": False, "detail": str(e), "changes": []}

@router.get("/figure-diagnostics")
async def get_figure_diagnostics(session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    input_path = state.get("content_pptx")
    if not input_path or not os.path.exists(input_path):
        return {"ok": False, "missing": [], "unplaced": []}
    try:
        session_upload_dir = get_session_upload_dir(session_id)
        extracts_dir = os.path.join(session_upload_dir, "pdf_extracts")
        missing, unplaced = collect_figure_diagnostics(input_path, extracts_dir)
        return {"ok": True, "missing": missing, "unplaced": unplaced}
    except Exception as e:
        return {"ok": False, "detail": str(e), "missing": [], "unplaced": []}

@router.get("/accessibility-report")
async def get_accessibility_report(session_id: str = Depends(get_session_id)):
    state = get_session_state(session_id)
    session_upload_dir = get_session_upload_dir(session_id)
    output_path = state.get("styled_pptx") or os.path.join(session_upload_dir, "styled_output.pptx")
    if not os.path.exists(output_path):
        return {"ok": False, "issues": []}
    try:
        issues = check_ppt_accessibility(output_path)
        return {"ok": True, "issues": issues}
    except Exception as e:
        return {"ok": False, "detail": str(e), "issues": []}

@router.get("/customers")
async def get_customers():
    try:
        json_path = os.path.join(os.path.dirname(__file__), "services", "customers.json")
        with open(json_path, "r") as f:
            data = json.load(f)
        return {"ok": True, "customers": data}
    except Exception as e:
        return {"ok": False, "detail": str(e), "customers": []}

@router.get("/media/{session_id}/{filename}")
async def serve_media(session_id: str, filename: str):
    clean_session_id = re.sub(r'[^a-zA-Z0-9_-]', '', session_id)
    clean_filename = re.sub(r'[^a-zA-Z0-9_.-]', '', filename)
    target = (Path(UPLOAD_DIR) / clean_session_id / clean_filename).resolve()
    if not target.is_file():
        target = (Path(TEMPLATES_DIR) / clean_filename).resolve()
    if not target.is_file():
        raise HTTPException(status_code=404, detail="File not found")
    return FileResponse(target)
