import os
import re
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from openpyxl import Workbook
from openpyxl.styles import Font, Alignment, PatternFill, Border, Side
from openpyxl.utils import get_column_letter

from .report import collect_changes, collect_figure_diagnostics
from .accessibility import check_ppt_accessibility

def create_excel_report(input_pptx, output_pptx, extracts_dir, customer_name, project_name, output_excel_path):
    wb = Workbook()
    
    navy_fill = PatternFill(start_color="1B365D", end_color="1B365D", fill_type="solid")
    cream_fill = PatternFill(start_color="F5F2EB", end_color="F5F2EB", fill_type="solid")
    green_fill = PatternFill(start_color="E2F0D9", end_color="E2F0D9", fill_type="solid")
    red_fill = PatternFill(start_color="FCE4D6", end_color="FCE4D6", fill_type="solid")
    header_gray_fill = PatternFill(start_color="E0E0E0", end_color="E0E0E0", fill_type="solid")
    
    font_title = Font(name="Calibri", size=16, bold=True, color="1B365D")
    font_header_white = Font(name="Calibri", size=11, bold=True, color="FFFFFF")
    font_header_dark = Font(name="Calibri", size=11, bold=True, color="1B365D")
    font_bold_data = Font(name="Calibri", size=11, bold=True)
    font_data = Font(name="Calibri", size=11)
    
    thin_border_side = Side(border_style="thin", color="D3D3D3")
    thin_border = Border(left=thin_border_side, right=thin_border_side, top=thin_border_side, bottom=thin_border_side)
    thick_bottom_side = Side(border_style="medium", color="1B365D")
    header_border = Border(left=thin_border_side, right=thin_border_side, top=thin_border_side, bottom=thick_bottom_side)

    prs = Presentation(output_pptx) if os.path.exists(output_pptx) else Presentation(input_pptx)
    total_slides = len(prs.slides)
    
    validation_checklist = []
    total_img_placeholders = 0
    empty_img_placeholders = 0
    
    for slide_idx, slide in enumerate(prs.slides):
        slide_num = slide_idx + 1
        missing_items = []
        
        for shape in slide.shapes:
            if shape.is_placeholder:
                ph_type = str(shape.placeholder_format.type)
                if "PICTURE" in ph_type or "BITMAP" in ph_type:
                    total_img_placeholders += 1
                    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
                        empty_img_placeholders += 1
                        missing_items.append("Picture Box")
            
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                ref_pattern = re.compile(
                    r"\b(figure|fig\.?|f\.?|table|tab\.?|t\.?|chart|image)\s*[\d.-]+|\binsert\s+(figure|fig\.?|f\.?|table|tab\.?|t\.?|chart|image)",
                    re.IGNORECASE
                )
                if ref_pattern.search(text):
                    missing_items.append(f"Missing image for reference: '{text}'")
                    
        status = "Complete" if not missing_items else "Action Required"
        details = ", ".join(missing_items) if missing_items else "All elements formatted"
        validation_checklist.append({
            "slide": slide_num,
            "status": status,
            "details": details
        })
        
    style_changes = []
    try:
        raw_changes = collect_changes(input_pptx, output_pptx)
        for s in raw_changes:
            slide_num = s.get("slide")
            for ph in s.get("placeholders", []):
                ph_name = ph.get("ph_type") or f"Placeholder {ph.get('ph_idx')}"
                for para in ph.get("paras", []):
                    para_idx = para.get("para_idx", 0) + 1
                    text_sample = para.get("text", "")
                    for chg in para.get("changes", []):
                        style_changes.append({
                            "slide": slide_num,
                            "ph_name": ph_name,
                            "para_idx": para_idx,
                            "text": text_sample,
                            "property": chg.get("prop"),
                            "val_in": str(chg.get("val_in")),
                            "val_out": str(chg.get("val_out"))
                        })
    except Exception as e:
        print("Excel Gen: Failed to collect style changes:", e)

    missing_figs = []
    unplaced_figs = []
    try:
        missing_figs, unplaced_figs = collect_figure_diagnostics(input_pptx, extracts_dir)
    except Exception as e:
        print("Excel Gen: Failed to collect figure diagnostics:", e)
        
    accessibility_issues = []
    try:
        accessibility_issues = check_ppt_accessibility(output_pptx)
    except Exception as e:
        print("Excel Gen: Failed to collect accessibility report:", e)

    ws_summary = wb.active
    ws_summary.title = "Summary Page"
    ws_summary.views.sheetView[0].showGridLines = True
    
    ws_summary["A1"] = "DECKFORGE COMPILATION SUMMARY REPORT"
    ws_summary["A1"].font = font_title
    ws_summary.row_dimensions[1].height = 30
    
    summary_data = [
        ("Customer Name", customer_name),
        ("Project Name", project_name),
        ("Input PPT Name", os.path.basename(input_pptx)),
        ("Total Slides", total_slides),
        ("Total Image Placeholders", total_img_placeholders),
        ("Insert Placeholders (Empty)", empty_img_placeholders),
    ]
    
    c1 = ws_summary.cell(row=3, column=1, value="Metric")
    c1.fill = navy_fill
    c1.font = font_header_white
    c1.border = header_border
    
    c2 = ws_summary.cell(row=3, column=2, value="Value")
    c2.fill = navy_fill
    c2.font = font_header_white
    c2.border = header_border
    
    for r_idx, (metric, val) in enumerate(summary_data, start=4):
        c1 = ws_summary.cell(row=r_idx, column=1, value=metric)
        c2 = ws_summary.cell(row=r_idx, column=2, value=val)
        c1.font = font_bold_data
        c2.font = font_data
        c1.border = thin_border
        c2.border = thin_border
        c1.fill = cream_fill
        
        if metric == "Insert Placeholders (Empty)":
            if val > 0:
                c2.fill = red_fill
                c2.font = Font(name="Calibri", size=11, color="9C0006", bold=True)
            else:
                c2.fill = green_fill
                c2.font = Font(name="Calibri", size=11, color="006100", bold=True)

    ws_val = wb.create_sheet(title="Slide Validation Check")
    ws_val.views.sheetView[0].showGridLines = True
    headers_val = ["Slide Number", "Status", "Missing Items / Details"]
    for col_idx, h in enumerate(headers_val, start=1):
        cell = ws_val.cell(row=1, column=col_idx, value=h)
        cell.fill = navy_fill
        cell.font = font_header_white
        cell.border = header_border
        cell.alignment = Alignment(horizontal="center" if col_idx < 3 else "left")
        
    for r_idx, item in enumerate(validation_checklist, start=2):
        c1 = ws_val.cell(row=r_idx, column=1, value=f"Slide {item['slide']}")
        c2 = ws_val.cell(row=r_idx, column=2, value=item["status"])
        c3 = ws_val.cell(row=r_idx, column=3, value=item["details"])
        
        for c in (c1, c2, c3):
            c.font = font_data
            c.border = thin_border
        
        c1.alignment = Alignment(horizontal="center")
        c2.alignment = Alignment(horizontal="center")
        
        if item["status"] == "Complete":
            c2.fill = green_fill
            c2.font = Font(name="Calibri", size=11, color="006100", bold=True)
        else:
            c2.fill = red_fill
            c2.font = Font(name="Calibri", size=11, color="9C0006", bold=True)

    ws_style = wb.create_sheet(title="Style Change Report")
    ws_style.views.sheetView[0].showGridLines = True
    headers_style = ["Slide Number", "Placeholder Name", "Paragraph Index", "Text Sample", "Property Changed", "Template Value (Target)", "Output Value (Result)"]
    for col_idx, h in enumerate(headers_style, start=1):
        cell = ws_style.cell(row=1, column=col_idx, value=h)
        cell.fill = navy_fill
        cell.font = font_header_white
        cell.border = header_border
        cell.alignment = Alignment(horizontal="center" if col_idx in (1, 3) else "left")
        
    for r_idx, chg in enumerate(style_changes, start=2):
        c1 = ws_style.cell(row=r_idx, column=1, value=f"Slide {chg['slide']}")
        c2 = ws_style.cell(row=r_idx, column=2, value=chg["ph_name"])
        c3 = ws_style.cell(row=r_idx, column=3, value=chg["para_idx"])
        c4 = ws_style.cell(row=r_idx, column=4, value=chg["text"])
        c5 = ws_style.cell(row=r_idx, column=5, value=chg["property"])
        c6 = ws_style.cell(row=r_idx, column=6, value=chg["val_in"])
        c7 = ws_style.cell(row=r_idx, column=7, value=chg["val_out"])
        
        for c in (c1, c2, c3, c4, c5, c6, c7):
            c.font = font_data
            c.border = thin_border
            
        c1.alignment = Alignment(horizontal="center")
        c3.alignment = Alignment(horizontal="center")

    ws_fig = wb.create_sheet(title="Figure Diagnostics")
    ws_fig.views.sheetView[0].showGridLines = True
    headers_fig = ["Category", "Figure/Crop Name", "Description"]
    for col_idx, h in enumerate(headers_fig, start=1):
        cell = ws_fig.cell(row=1, column=col_idx, value=h)
        cell.fill = navy_fill
        cell.font = font_header_white
        cell.border = header_border
        cell.alignment = Alignment(horizontal="center" if col_idx == 1 else "left")
        
    fig_rows = []
    for fig in missing_figs:
        fig_rows.append(("Missing Figure (Requested in PPT but not found in PDF Extracts)", fig, f"The presentation requests '{fig}' but no matching cropped figure image was uploaded."))
    for fig in unplaced_figs:
        fig_rows.append(("Unplaced Crop (Uploaded but not placed in PPT)", fig, f"The figure '{fig}' was cropped/uploaded, but no corresponding placeholder matches it in the presentation."))
        
    for r_idx, (cat, name, desc) in enumerate(fig_rows, start=2):
        c1 = ws_fig.cell(row=r_idx, column=1, value=cat)
        c2 = ws_fig.cell(row=r_idx, column=2, value=name)
        c3 = ws_fig.cell(row=r_idx, column=3, value=desc)
        
        for c in (c1, c2, c3):
            c.font = font_data
            c.border = thin_border
            
        if "Missing" in cat:
            c1.fill = red_fill
            c1.font = Font(name="Calibri", size=11, color="9C0006", bold=True)
        else:
            c1.fill = cream_fill
            c1.font = Font(name="Calibri", size=11, color="595959", bold=True)

    ws_acc = wb.create_sheet(title="Accessibility Report")
    ws_acc.views.sheetView[0].showGridLines = True
    headers_acc = ["Slide Number", "Category", "Severity", "Detail Description"]
    for col_idx, h in enumerate(headers_acc, start=1):
        cell = ws_acc.cell(row=1, column=col_idx, value=h)
        cell.fill = navy_fill
        cell.font = font_header_white
        cell.border = header_border
        cell.alignment = Alignment(horizontal="center" if col_idx in (1, 3) else "left")
        
    for r_idx, issue in enumerate(accessibility_issues, start=2):
        slide_val = f"Slide {issue.get('slide', 0) + 1}" if issue.get("slide") is not None else "General"
        c1 = ws_acc.cell(row=r_idx, column=1, value=slide_val)
        c2 = ws_acc.cell(row=r_idx, column=2, value=issue.get("category"))
        c3 = ws_acc.cell(row=r_idx, column=3, value=issue.get("severity"))
        c4 = ws_acc.cell(row=r_idx, column=4, value=issue.get("detail"))
        
        for c in (c1, c2, c3, c4):
            c.font = font_data
            c.border = thin_border
            
        c1.alignment = Alignment(horizontal="center")
        c3.alignment = Alignment(horizontal="center")
        
        sev = issue.get("severity", "").upper()
        if "ERROR" in sev:
            c3.fill = red_fill
            c3.font = Font(name="Calibri", size=11, color="9C0006", bold=True)
        elif "WARNING" in sev:
            c3.fill = PatternFill(start_color="FFF2CC", end_color="FFF2CC", fill_type="solid")
            c3.font = Font(name="Calibri", size=11, color="7F6000", bold=True)

    for ws in wb.worksheets:
        for col in ws.columns:
            max_len = 0
            for cell in col:
                val = str(cell.value or '')
                if '\n' in val:
                    val = max(val.split('\n'), key=len)
                if len(val) > max_len:
                    max_len = len(val)
            col_letter = get_column_letter(col[0].column)
            ws.column_dimensions[col_letter].width = min(max(max_len + 3, 12), 65)

    wb.save(output_excel_path)
