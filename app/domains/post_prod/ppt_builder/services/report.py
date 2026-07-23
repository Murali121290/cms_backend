from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from lxml import etree
import json, sys, os, re

INPUT_PATH  = "input.pptx"
OUTPUT_PATH = "output.pptx"
REPORT_PATH = "change_report.html"

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS   = {"a": A_NS, "p": P_NS}

ALIGN_NAMES = {
    PP_ALIGN.LEFT: "Left", PP_ALIGN.CENTER: "Center",
    PP_ALIGN.RIGHT: "Right", PP_ALIGN.JUSTIFY: "Justify",
    PP_ALIGN.DISTRIBUTE: "Distribute",
}
SCHEME_MAP = {
    "tx1": "dk1", "bg1": "lt1", "tx2": "dk2", "bg2": "lt2",
    "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
    "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
    "hlink": "hlink", "folHlink": "folHlink",
}
BODY_LEVELS = ["lvl1pPr","lvl2pPr","lvl3pPr","lvl4pPr","lvl5pPr",
               "lvl6pPr","lvl7pPr","lvl8pPr","lvl9pPr"]
TITLE_TYPES = {"TITLE (1)", "CENTER_TITLE (3)"}
THEME_FONTS = {"+mn-lt","+mj-lt","+mn-ea","+mj-ea","+mn-cs","+mj-cs"}

def get_theme_colors(prs):
    master_part = prs.slide_master.part
    for rel in master_part.rels.values():
        if "theme" in rel.reltype:
            try:
                theme_el = etree.fromstring(rel.target_part.blob)
                colors = {}
                clr = theme_el.find(".//a:clrScheme", NS)
                if clr is not None:
                    for child in clr:
                        name = child.tag.split("}")[-1]
                        srgb   = child.find(".//a:srgbClr", NS)
                        sysClr = child.find(".//a:sysClr",  NS)
                        if srgb is not None:
                            colors[name] = "#" + srgb.get("val", "").upper()
                        elif sysClr is not None:
                            v = sysClr.get("lastClr") or sysClr.get("val", "")
                            if v: colors[name] = "#" + v.upper()
                return colors
            except Exception:
                pass
    return {}

def resolve_color_el(el, colors):
    if el is None:
        return None
    srgb   = el.find(".//a:srgbClr", NS)
    scheme = el.find(".//a:schemeClr", NS)
    sysClr = el.find(".//a:sysClr",   NS)
    if srgb is not None:
        return "#" + srgb.get("val", "").upper()
    if scheme is not None:
        token = scheme.get("val", "")
        key   = SCHEME_MAP.get(token, token)
        return colors.get(key)
    if sysClr is not None:
        v = sysClr.get("lastClr") or sysClr.get("val", "")
        return ("#" + v.upper()) if v else None
    return None

def parse_level_style(lvl_el, colors):
    s = {}
    algn = lvl_el.get("algn")
    if algn:
        s["alignment"] = {"l":"Left","ctr":"Center","r":"Right",
                          "just":"Justify","dist":"Distribute"}.get(algn, algn)
    defRPr = lvl_el.find("a:defRPr", NS)
    if defRPr is not None:
        sz = defRPr.get("sz")
        if sz:
            s["fontSize_pt"] = str(round(int(sz) / 100, 1)) + "pt"
        b = defRPr.get("b")
        if b is not None:
            s["bold"] = "Yes" if b == "1" else "No"
        i = defRPr.get("i")
        if i is not None:
            s["italic"] = "Yes" if i == "1" else "No"
        latin = defRPr.find("a:latin", NS)
        if latin is not None:
            f = latin.get("typeface", "")
            if f and f not in THEME_FONTS:
                s["fontFamily"] = f
        fill = defRPr.find("a:solidFill", NS)
        clr = resolve_color_el(fill, colors)
        if clr:
            s["color"] = clr
    return s

def get_master_styles(prs, colors):
    txStyles = prs.slide_master.element.find("p:txStyles", NS)
    if txStyles is None:
        return {}
    out = {}
    for style_el in txStyles:
        name = style_el.tag.split("}")[-1]
        levels = {}
        for lvl_el in style_el:
            tag   = lvl_el.tag.split("}")[-1]
            style = parse_level_style(lvl_el, colors)
            if style:
                levels[tag] = style
        if levels:
            out[name] = levels
    return out

def effective_style(master_styles, ph_type_str, level):
    style_name = "titleStyle" if any(t in ph_type_str for t in TITLE_TYPES) else "bodyStyle"
    level_key  = BODY_LEVELS[min(level, len(BODY_LEVELS) - 1)]
    lvl = master_styles.get(style_name, {}).get(level_key, {})
    if not lvl:
        lvl = master_styles.get(style_name, {}).get("lvl1pPr", {})
    return lvl

def pt_str(size):
    try:
        return str(round(size.pt, 1)) + "pt" if size else None
    except Exception:
        return None

def rgb_str(font):
    try:
        if font.color and font.color.type is not None:
            return "#" + str(font.color.rgb).upper()
    except Exception:
        pass
    return None

def lnspc_str(para):
    v = para.line_spacing
    if v is None:
        return None
    if isinstance(v, float):
        return f"{round(v * 100)}%"
    try:
        return f"{v.pt}pt"
    except Exception:
        return str(v)

def align_str(para):
    return ALIGN_NAMES.get(para.alignment)

def collect_changes(input_path, output_path):
    prs_in  = Presentation(input_path)
    prs_out = Presentation(output_path)

    colors_in      = get_theme_colors(prs_in)
    master_styles  = get_master_styles(prs_in, colors_in)

    slides_data = []

    for si, (sl_in, sl_out) in enumerate(zip(prs_in.slides, prs_out.slides)):
        shapes_in  = {s.shape_id: s for s in sl_in.shapes
                      if s.is_placeholder and s.has_text_frame}
        shapes_out = {s.shape_id: s for s in sl_out.shapes
                      if s.is_placeholder and s.has_text_frame}

        placeholders = []

        for sid in sorted(set(shapes_in) & set(shapes_out)):
            sh_in, sh_out = shapes_in[sid], shapes_out[sid]
            ph_type = str(sh_in.placeholder_format.type)

            paras_changed = []

            for pi, (p_in, p_out) in enumerate(
                    zip(sh_in.text_frame.paragraphs,
                        sh_out.text_frame.paragraphs)):

                level = p_in.level or 0
                eff   = effective_style(master_styles, ph_type, level)
                changes = []

                a_in  = align_str(p_in)  or eff.get("alignment")
                a_out = align_str(p_out) or eff.get("alignment")
                if a_in != a_out:
                    changes.append({"prop":"Alignment",
                                    "before": a_in or "—", "after": a_out or "—"})

                ls_in  = lnspc_str(p_in)
                ls_out = lnspc_str(p_out)
                if ls_in != ls_out and ls_out:
                    changes.append({"prop":"Line spacing",
                                    "before": ls_in or "—", "after": ls_out})

                sb_in  = pt_str(p_in.space_before)
                sb_out = pt_str(p_out.space_before)
                if sb_in != sb_out and sb_out:
                    changes.append({"prop":"Space before",
                                    "before": sb_in or "—", "after": sb_out})

                for ri, (r_in, r_out) in enumerate(zip(p_in.runs, p_out.runs)):
                    fi, fo = r_in.font, r_out.font

                    fs_in  = pt_str(fi.size) or eff.get("fontSize_pt")
                    fs_out = pt_str(fo.size)
                    if fs_in != fs_out and fs_out:
                        changes.append({"prop":"Font size",
                                        "before": fs_in or "—", "after": fs_out})

                    ff_in  = fi.name if fi.name and fi.name not in THEME_FONTS else eff.get("fontFamily")
                    ff_out = fo.name if fo.name and fo.name not in THEME_FONTS else None
                    if ff_in != ff_out and ff_out:
                        changes.append({"prop":"Font family",
                                        "before": ff_in or "—", "after": ff_out})

                    b_in  = ("Yes" if fi.bold else "No") if fi.bold is not None else eff.get("bold")
                    b_out = ("Yes" if fo.bold else "No") if fo.bold is not None else None
                    if b_in != b_out and b_out:
                        changes.append({"prop":"Bold",
                                        "before": b_in or "—", "after": b_out})

                    c_in  = rgb_str(fi) or eff.get("color")
                    c_out = rgb_str(fo)
                    if c_in != c_out and c_out:
                        changes.append({"prop":"Color", "is_color": True,
                                        "before": c_in or "—", "after": c_out})

                if changes:
                    paras_changed.append({
                        "text":    p_in.text[:55] if p_in.text else "",
                        "idx":     pi,
                        "changes": changes,
                    })

            if paras_changed:
                placeholders.append({
                    "name":  sh_in.name,
                    "type":  ph_type,
                    "idx":   sh_in.placeholder_format.idx,
                    "paras": paras_changed,
                })

        if placeholders:
            slides_data.append({"slide": si + 1, "placeholders": placeholders})

    return slides_data

def collect_figure_diagnostics(input_path, extracts_dir):
    if not os.path.exists(input_path):
        return [], []

    requested_figs = set()
    try:
        prs = Presentation(input_path)
        _FIG_PAT = re.compile(r'insert\s+(figure|table)\s+([\d.-]+)', re.IGNORECASE)
        for slide in prs.slides:
            all_shapes = []
            def recurse(container):
                for sh in container:
                    if sh.shape_type == 6:
                        try:
                            recurse(sh.shapes)
                        except Exception:
                            pass
                    else:
                        all_shapes.append(sh)
            recurse(slide.shapes)
            
            for shape in all_shapes:
                if shape.has_text_frame:
                    for m in _FIG_PAT.finditer(shape.text_frame.text):
                        requested_figs.add(f"{m.group(1).lower()} {m.group(2)}")
    except Exception as e:
        print("Error checking requested figures:", e)

    cropped_figs = set()
    if os.path.isdir(extracts_dir):
        for fname in os.listdir(extracts_dir):
            if fname.lower().endswith(('.png', '.jpg', '.jpeg')):
                key = os.path.splitext(fname)[0].strip().lower()
                cropped_figs.add(key)

    missing_figs = sorted(list(requested_figs - cropped_figs))
    unplaced_figs = sorted(list(cropped_figs - requested_figs))
    return missing_figs, unplaced_figs
