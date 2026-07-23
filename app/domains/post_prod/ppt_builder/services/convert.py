from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree
import copy
import json
import re
import sys
import os
import io

INPUT_PATH = "input.pptx"
TEMPLATE_STYLE_PATH = "template_styles.json"
OUTPUT_PATH = "output.pptx"
COVER_IMAGE_PATH = "Picture1.png"

PT_TO_EMU = 12700
A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"

ALIGN_MAP = {
    "l":    PP_ALIGN.LEFT,
    "ctr":  PP_ALIGN.CENTER,
    "r":    PP_ALIGN.RIGHT,
    "just": PP_ALIGN.JUSTIFY,
    "dist": PP_ALIGN.DISTRIBUTE,
}

SCHEME_TOKEN_MAP = {
    "tx1": "dk1",  "bg1": "lt1",
    "tx2": "dk2",  "bg2": "lt2",
    "accent1": "accent1", "accent2": "accent2", "accent3": "accent3",
    "accent4": "accent4", "accent5": "accent5", "accent6": "accent6",
    "hlink": "hlink",     "folHlink": "folHlink",
}

THEME_FONTS = {"+mn-lt", "+mj-lt", "+mn-ea", "+mj-ea", "+mn-cs", "+mj-cs"}
TITLE_PH_TYPES = {"TITLE (1)", "CENTER_TITLE (3)"}

BODY_LEVEL_KEYS = [
    "lvl1pPr", "lvl2pPr", "lvl3pPr", "lvl4pPr", "lvl5pPr",
    "lvl6pPr", "lvl7pPr", "lvl8pPr", "lvl9pPr",
]

THEME_FONT_KEY = {
    "+mj-lt": "major", "+mj-ea": "major", "+mj-cs": "major",
    "+mn-lt": "minor", "+mn-ea": "minor", "+mn-cs": "minor",
}

def resolve_color(color_str, color_scheme):
    if not color_str or color_str == "none":
        return None
    if color_str.startswith("#"):
        h = color_str[1:]
        return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    if color_str.startswith("scheme:"):
        token = color_str[7:]
        key = SCHEME_TOKEN_MAP.get(token, token)
        h = color_scheme.get(key, "")
        if h and h.startswith("#"):
            h = h[1:]
            return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
    return None

def _get_shape_spPr(shape):
    spPr = shape._element.find(f"{{{P_NS}}}spPr")
    if spPr is None:
        spPr = shape._element.find(f"{{{A_NS}}}spPr")
    return spPr

def _get_or_create_pPr(p_el):
    pPr = p_el.find(f"{{{A_NS}}}pPr")
    if pPr is None:
        pPr = etree.Element(f"{{{A_NS}}}pPr")
        p_el.insert(0, pPr)
    return pPr

def _clear_fill(parent_el):
    for tag in [f"{{{A_NS}}}solidFill", f"{{{A_NS}}}noFill",
                f"{{{A_NS}}}gradFill", f"{{{A_NS}}}pattFill", f"{{{A_NS}}}blipFill"]:
        for child in parent_el.findall(tag):
            parent_el.remove(child)

def _set_solid_fill(parent_el, rgb):
    _clear_fill(parent_el)
    solid = etree.SubElement(parent_el, f"{{{A_NS}}}solidFill")
    srgb = etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
    srgb.set("val", str(rgb))

def apply_run_style(run, style, color_scheme, font_scheme=None):
    if not style:
        return
    font = run.font
    rPr = font._element

    if "fontSize_pt" in style:
        font.size = Pt(style["fontSize_pt"])

    if "fontFamily" in style:
        fn = style["fontFamily"]
        if fn in THEME_FONTS:
            key = THEME_FONT_KEY.get(fn, "minor")
            fn = (font_scheme or {}).get(key) if font_scheme else None
        if fn:
            font.name = fn

    if "bold" in style:
        font.bold = style["bold"]
    if "italic" in style:
        font.italic = style["italic"]
    if "underline" in style and isinstance(style["underline"], bool):
        font.underline = style["underline"]
    if "color" in style and (font.color is None or font.color.type is None):
        rgb = resolve_color(style["color"], color_scheme)
        if rgb:
            try:
                font.color.rgb = rgb
            except Exception:
                pass

    if rPr is not None:
        if "kerning_pt" in style:
            rPr.set("kern", str(int(style["kerning_pt"] * 100)))
        if "charSpacing_pt" in style:
            rPr.set("spc", str(int(style["charSpacing_pt"] * 100)))
        if "language" in style:
            rPr.set("lang", style["language"])

def apply_para_style(para, style, color_scheme):
    if not style:
        return

    if "alignment" in style:
        pPr_existing = para._p.find(f"{{{A_NS}}}pPr")
        has_explicit_algn = (pPr_existing is not None and
                             pPr_existing.get("algn") is not None)
        if not has_explicit_algn:
            para.alignment = ALIGN_MAP.get(style["alignment"])

    if "spaceBefore_pt" in style:
        para.space_before = Pt(style["spaceBefore_pt"])
    if "spaceAfter_pt" in style:
        para.space_after = Pt(style["spaceAfter_pt"])

    if "lineSpacing_pct" in style:
        para.line_spacing = style["lineSpacing_pct"] / 100.0
    elif "lineSpacing_pt" in style:
        para.line_spacing = Pt(style["lineSpacing_pt"])

    p_el    = para._p
    pPr_el  = p_el.find(f"{{{A_NS}}}pPr")
    has_bu_none = (pPr_el is not None and
                   pPr_el.find(f"{{{A_NS}}}buNone") is not None)

    if not has_bu_none:
        if "marginLeft_pt" in style or "indent_pt" in style:
            pPr = _get_or_create_pPr(p_el)
            if "marginLeft_pt" in style and pPr.get("marL") is None:
                pPr.set("marL", str(int(style["marginLeft_pt"] * PT_TO_EMU)))
            if "indent_pt" in style and pPr.get("indent") is None:
                pPr.set("indent", str(int(style["indent_pt"] * PT_TO_EMU)))

        if "bulletChar" in style or "bulletFont" in style:
            pPr = _get_or_create_pPr(p_el)
            has_auto_num = pPr.find(f"{{{A_NS}}}buAutoNum") is not None
            
            if has_auto_num:
                for bu_tag in [f"{{{A_NS}}}buFontTx", f"{{{A_NS}}}buChar", f"{{{A_NS}}}buBlip"]:
                    for el in pPr.findall(bu_tag):
                        pPr.remove(el)
            else:
                for bu_tag in [f"{{{A_NS}}}buFontTx", f"{{{A_NS}}}buFont",
                               f"{{{A_NS}}}buChar", f"{{{A_NS}}}buAutoNum", f"{{{A_NS}}}buBlip"]:
                    for el in pPr.findall(bu_tag):
                        pPr.remove(el)
                if "bulletFont" in style:
                    bu_font = etree.SubElement(pPr, f"{{{A_NS}}}buFont")
                    bu_font.set("typeface", style["bulletFont"])
                if "bulletChar" in style:
                    bu_char = etree.SubElement(pPr, f"{{{A_NS}}}buChar")
                    bu_char.set("char", style["bulletChar"])

def apply_body_properties(tf, body_props):
    txBody = tf._txBody
    bodyPr = txBody.find(f"{{{A_NS}}}bodyPr")
    if bodyPr is None:
        return

    if body_props:
        for key, attr in [("verticalAlignment", "anchor"), ("textDirection", "vert"), ("wrap", "wrap")]:
            if key in body_props:
                bodyPr.set(attr, str(body_props[key]))
        for key, attr in [("lIns_pt", "lIns"), ("rIns_pt", "rIns"), ("tIns_pt", "tIns"), ("bIns_pt", "bIns")]:
            if key in body_props:
                bodyPr.set(attr, str(int(body_props[key] * PT_TO_EMU)))

    for af_tag in [f"{{{A_NS}}}normAutofit", f"{{{A_NS}}}spAutoFit", f"{{{A_NS}}}noAutofit"]:
        for old in bodyPr.findall(af_tag):
            bodyPr.remove(old)
    etree.SubElement(bodyPr, f"{{{A_NS}}}normAutofit")

def apply_shape_geometry(shape, template_shape):
    if not template_shape:
        return
    try:
        if "position" in template_shape:
            pos = template_shape["position"]
            if pos.get("x_pt") is not None:
                shape.left = int(pos["x_pt"] * PT_TO_EMU)
            if pos.get("y_pt") is not None:
                shape.top = int(pos["y_pt"] * PT_TO_EMU)
        if "size" in template_shape:
            sz = template_shape["size"]
            if sz.get("width_pt") is not None:
                shape.width = int(sz["width_pt"] * PT_TO_EMU)
            if sz.get("height_pt") is not None:
                shape.height = int(sz["height_pt"] * PT_TO_EMU)
    except Exception:
        pass

def apply_shape_fill(shape, template_shape, color_scheme):
    if not template_shape or "fill" not in template_shape:
        return
    fill_val = template_shape["fill"]
    spPr = _get_shape_spPr(shape)
    if spPr is None:
        return
    _clear_fill(spPr)
    if fill_val == "none":
        etree.SubElement(spPr, f"{{{A_NS}}}noFill")
    else:
        rgb = resolve_color(fill_val, color_scheme)
        if rgb:
            _set_solid_fill(spPr, rgb)

def apply_shape_border(shape, template_shape, color_scheme):
    if not template_shape or "border" not in template_shape:
        return
    border = template_shape["border"]
    spPr = _get_shape_spPr(shape)
    if spPr is None:
        return
    for child in spPr.findall(f"{{{A_NS}}}ln"):
        spPr.remove(child)
    ln = etree.SubElement(spPr, f"{{{A_NS}}}ln")
    if "width_pt" in border:
        ln.set("w", str(int(border["width_pt"] * PT_TO_EMU)))
    color = border.get("color")
    if color == "none":
        etree.SubElement(ln, f"{{{A_NS}}}noFill")
    elif color:
        rgb = resolve_color(color, color_scheme)
        if rgb:
            solid = etree.SubElement(ln, f"{{{A_NS}}}solidFill")
            srgb = etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
            srgb.set("val", str(rgb))
    if "dashStyle" in border:
        dash = etree.SubElement(ln, f"{{{A_NS}}}prstDash")
        dash.set("val", border["dashStyle"])

def _make_decorative_sp(sh_data, shape_id, color_scheme):
    sp = etree.Element(f"{{{P_NS}}}sp")
    nvSpPr = etree.SubElement(sp, f"{{{P_NS}}}nvSpPr")
    cNvPr = etree.SubElement(nvSpPr, f"{{{P_NS}}}cNvPr")
    cNvPr.set("id", str(shape_id))
    cNvPr.set("name", sh_data.get("shapeName", f"deco{shape_id}"))
    etree.SubElement(nvSpPr, f"{{{P_NS}}}cNvSpPr")
    etree.SubElement(nvSpPr, f"{{{P_NS}}}nvPr")

    spPr = etree.SubElement(sp, f"{{{P_NS}}}spPr")
    pos  = sh_data["position"]
    sz   = sh_data["size"]
    xfrm = etree.SubElement(spPr, f"{{{A_NS}}}xfrm")
    off  = etree.SubElement(xfrm, f"{{{A_NS}}}off")
    off.set("x", str(int(pos["x_pt"] * PT_TO_EMU)))
    off.set("y", str(int(pos["y_pt"] * PT_TO_EMU)))
    ext  = etree.SubElement(xfrm, f"{{{A_NS}}}ext")
    ext.set("cx", str(int(sz["width_pt"] * PT_TO_EMU)))
    ext.set("cy", str(int(sz["height_pt"] * PT_TO_EMU)))

    prstGeom = etree.SubElement(spPr, f"{{{A_NS}}}prstGeom")
    prstGeom.set("prst", "rect")
    etree.SubElement(prstGeom, f"{{{A_NS}}}avLst")

    fill_val = sh_data.get("fill")
    if fill_val == "none":
        etree.SubElement(spPr, f"{{{A_NS}}}noFill")
    elif fill_val:
        rgb = resolve_color(fill_val, color_scheme)
        if rgb:
            solid = etree.SubElement(spPr, f"{{{A_NS}}}solidFill")
            srgb  = etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
            srgb.set("val", str(rgb))

    border = sh_data.get("border", {})
    border_color = border.get("color")
    if border_color == "none":
        ln = etree.SubElement(spPr, f"{{{A_NS}}}ln")
        etree.SubElement(ln, f"{{{A_NS}}}noFill")
    elif border_color:
        rgb = resolve_color(border_color, color_scheme)
        ln  = etree.SubElement(spPr, f"{{{A_NS}}}ln")
        if "width_pt" in border:
            ln.set("w", str(int(border["width_pt"] * PT_TO_EMU)))
        if rgb:
            solid = etree.SubElement(ln, f"{{{A_NS}}}solidFill")
            srgb  = etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
            srgb.set("val", str(rgb))

    return sp

def _remove_copyright_from_spTree(spTree):
    for child in list(spTree):
        if child.find(f".//{{{P_NS}}}ph") is not None:
            continue
        c_text = "".join(t.text or "" for t in child.iter(f"{{{A_NS}}}t")).strip()
        if "©" in c_text:
            spTree.remove(child)
            continue
        if child.tag.split("}")[-1] == "grpSp":
            inner_text = "".join(t.text or "" for t in child.iter(f"{{{A_NS}}}t")).strip()
            if "©" in inner_text:
                spTree.remove(child)

def strip_copyright_from_masters(prs):
    for master in prs.slide_masters:
        master_spTree = master._element.find(f".//{{{P_NS}}}spTree")
        if master_spTree is not None:
            _remove_copyright_from_spTree(master_spTree)
        for layout in master.slide_layouts:
            layout_spTree = layout._element.find(f".//{{{P_NS}}}spTree")
            if layout_spTree is not None:
                _remove_copyright_from_spTree(layout_spTree)

def _relink_images(sp, source_shape, target_slide):
    source_part = source_shape.part
    for attr in (f"{{{R_NS}}}embed", f"{{{R_NS}}}link"):
        for blip in sp.iter(f"{{{A_NS}}}blip"):
            old_rid = blip.get(attr)
            if not old_rid:
                continue
            image_part = source_part.related_part(old_rid)
            _, new_rid = target_slide.part.get_or_add_image_part(io.BytesIO(image_part.blob))
            blip.set(attr, new_rid)

def add_decorative_shapes(slide, layout_name, template_json, color_scheme, template_prs=None):
    cSld   = slide._element.find(f"{{{P_NS}}}cSld")
    spTree = cSld.find(f"{{{P_NS}}}spTree") if cSld is not None else None
    if spTree is None:
        return

    if template_prs is not None:
        master = template_prs.slide_master
        try:
            layout_idx = int(layout_name.replace("slideLayout", "")) - 1
            layout_deco = [s for s in master.slide_layouts[layout_idx].shapes
                           if not s.is_placeholder]
        except (ValueError, IndexError):
            layout_deco = []

        master_deco = [s for s in master.shapes if not s.is_placeholder]

        _remove_copyright_from_spTree(spTree)

        insert_idx = 0
        for i, child in enumerate(spTree):
            if child.tag.split("}")[-1] in ("nvGrpSpPr", "grpSpPr"):
                insert_idx = i + 1

        next_id = max(
            (int(el.get("id")) for el in spTree.iter(f"{{{P_NS}}}cNvPr") if el.get("id", "").isdigit()),
            default=0,
        ) + 1

        dk1 = color_scheme.get("dk1", "")

        for sh in master_deco + layout_deco:
            sp = copy.deepcopy(sh._element)
            _relink_images(sp, sh, slide)
            cNvPr = sp.find(f".//{{{P_NS}}}cNvPr")
            if cNvPr is not None:
                cNvPr.set("id", str(next_id))
            if dk1.startswith("#"):
                hex_val = dk1[1:]
                for rPr in sp.iter(f"{{{A_NS}}}rPr"):
                    if rPr.find(f"{{{A_NS}}}solidFill") is None and \
                       rPr.find(f"{{{A_NS}}}schemeClr") is None:
                        solid = etree.Element(f"{{{A_NS}}}solidFill")
                        srgb  = etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
                        srgb.set("val", hex_val)
                        rPr.insert(0, solid)
            spTree.insert(insert_idx, sp)
            insert_idx += 1
            next_id    += 1
    else:
        sources = [sh for sh in template_json.get("masterShapes", []) if not sh.get("placeholder")]
        for layout in template_json["slideLayouts"]:
            if layout["layoutName"] == layout_name:
                sources += layout.get("decorativeShapes", [])
                break
        sources = [s for s in sources if s.get("fill") and s["fill"] != "none"
                   and s.get("position") and s.get("size")]

        _remove_copyright_from_spTree(spTree)

        insert_idx = 0
        for i, child in enumerate(spTree):
            if child.tag.split("}")[-1] in ("nvGrpSpPr", "grpSpPr"):
                insert_idx = i + 1

        next_id = max(
            (int(el.get("id")) for el in spTree.iter(f"{{{P_NS}}}cNvPr") if el.get("id", "").isdigit()),
            default=0,
        ) + 1

        for sh_data in sources:
            sp = _make_decorative_sp(sh_data, next_id, color_scheme)
            spTree.insert(insert_idx, sp)
            insert_idx += 1
            next_id    += 1

def apply_slide_background(slide, bg_fill, color_scheme):
    if not bg_fill or bg_fill == "none":
        return
    rgb = resolve_color(bg_fill, color_scheme)
    if not rgb:
        return
    cSld = slide._element.find(f"{{{P_NS}}}cSld")
    if cSld is None:
        return
    for bg_el in cSld.findall(f"{{{P_NS}}}bg"):
        cSld.remove(bg_el)
    bg_el = etree.Element(f"{{{P_NS}}}bg")
    bgPr = etree.SubElement(bg_el, f"{{{P_NS}}}bgPr")
    solid = etree.SubElement(bgPr, f"{{{A_NS}}}solidFill")
    srgb = etree.SubElement(solid, f"{{{A_NS}}}srgbClr")
    srgb.set("val", str(rgb))
    etree.SubElement(bgPr, f"{{{A_NS}}}effectLst")
    spTree = cSld.find(f"{{{P_NS}}}spTree")
    if spTree is not None:
        cSld.insert(list(cSld).index(spTree), bg_el)
    else:
        cSld.insert(0, bg_el)

def flatten_level_style(lvl_dict):
    style = {k: v for k, v in lvl_dict.items() if k != "defaultRunProps"}
    style.update(lvl_dict.get("defaultRunProps", {}))
    return style

def get_master_style_for_title(master_styles):
    lvl = master_styles.get("titleStyle", {}).get("lvl1pPr", {})
    return flatten_level_style(lvl)

def get_master_style_for_body(master_styles, level=0):
    level_key = BODY_LEVEL_KEYS[min(level, len(BODY_LEVEL_KEYS) - 1)]
    lvl = master_styles.get("bodyStyle", {}).get(level_key, {})
    return flatten_level_style(lvl)

def get_master_shape(template_json, ph_idx):
    for sh in template_json.get("masterShapes", []):
        if sh.get("placeholder", {}).get("idx") == ph_idx:
            return sh
    return None

def get_layout_ph(template_json, layout_name, ph_idx):
    for layout in template_json["slideLayouts"]:
        if layout["layoutName"] == layout_name:
            for ph in layout["placeholders"]:
                if ph.get("placeholder", {}).get("idx") == ph_idx:
                    return ph
    return None

def get_layout_style(template_json, layout_name, ph_idx, level=0):
    level_key = BODY_LEVEL_KEYS[min(level, len(BODY_LEVEL_KEYS) - 1)]
    for layout in template_json["slideLayouts"]:
        if layout["layoutName"] == layout_name:
            for ph in layout["placeholders"]:
                if ph.get("placeholder", {}).get("idx") == ph_idx:
                    tb = ph.get("textBody", {})
                    list_style = tb.get("listStyle", {})
                    lvl = list_style.get(level_key) or list_style.get("lvl1pPr", {})
                    return flatten_level_style(lvl)
    return {}

def get_template_slide_ph(template_json, slide_num, ph_idx):
    for slide in template_json["slides"]:
        if slide["slideNumber"] == slide_num:
            for sh in slide["shapes"]:
                if sh.get("placeholder", {}).get("idx") == ph_idx:
                    style = {}
                    for para in sh.get("textBody", {}).get("paragraphs", []):
                        para_style = {k: v for k, v in para.items()
                                      if k not in ("runs", "sampleText")}
                        runs = para.get("runs", [])
                        if runs:
                            run_style = {k: v for k, v in runs[0].items()
                                         if k != "sampleText"}
                            style.update(para_style)
                            style.update(run_style)
                            break
                    return style, sh
    return {}, None

def merge(*dicts):
    out = {}
    for d in dicts:
        if d:
            out.update(d)
    return out

def ph_type_category(ph_type_str):
    s = ph_type_str.upper()
    if "CENTER_TITLE" in s:
        return "center_title"
    if "SUBTITLE" in s:
        return "subtitle"
    if "TITLE" in s:
        return "title"
    if "PICTURE" in s:
        return "picture"
    if "BODY" in s or "OBJECT" in s or "TEXT" in s:
        return "body"
    return "other"

def input_slide_signature(slide):
    sig = {}
    nonph_text_count = 0
    has_picture = False
    for shape in slide.shapes:
        if shape.is_placeholder and shape.placeholder_format is not None:
            ph = shape.placeholder_format
            sig[ph.idx] = ph_type_category(str(ph.type))
        else:
            if shape.shape_type == 13:
                has_picture = True
            elif shape.has_text_frame:
                text = shape.text_frame.text.strip()
                if any(p.text.strip() for p in shape.text_frame.paragraphs):
                    if _FIG_PAT.search(text):
                        has_picture = True
                    else:
                        nonph_text_count += 1
    if has_picture:
        sig[10] = "picture"
    if 1 in sig and sig.get(1) == "body" and nonph_text_count > 0:
        next_idx = max(sig.keys()) + 1
        for _ in range(nonph_text_count):
            sig[next_idx] = "body"
            next_idx += 1
    return sig

def layout_signature(layout_phs):
    sig = {}
    for ph in layout_phs:
        ph_info = ph.get("placeholder", {})
        idx = ph_info.get("idx")
        cat = ph_type_category(ph_info.get("type", ""))
        if idx is not None:
            sig[idx] = cat
    return sig

def match_score(input_sig, template_sig):
    score = 0
    has_input_pic = any(cat == "picture" for cat in input_sig.values())
    has_tmpl_pic = any(cat == "picture" for cat in template_sig.values())
    if has_input_pic and has_tmpl_pic:
        score += 3
    for idx, cat in input_sig.items():
        if idx in template_sig:
            score += 2 if template_sig[idx] == cat else 1
    return score

def find_best_layout(input_slide, template_json):
    input_sig = input_slide_signature(input_slide)
    best_name, best_score = None, -1
    for layout in template_json["slideLayouts"]:
        t_sig = layout_signature(layout["placeholders"])
        score = match_score(input_sig, t_sig)
        if score > best_score:
            best_score, best_name = score, layout["layoutName"]
    return best_name, best_score

def find_best_template_slide(input_slide, template_json):
    input_sig = input_slide_signature(input_slide)
    best_num, best_score = None, -1
    for tslide in template_json["slides"]:
        t_sig = {}
        for sh in tslide["shapes"]:
            ph_info = sh.get("placeholder", {})
            idx = ph_info.get("idx")
            cat = ph_type_category(ph_info.get("type", ""))
            if idx is not None:
                t_sig[idx] = cat
        score = match_score(input_sig, t_sig)
        if score > best_score:
            best_score, best_num = score, tslide["slideNumber"]
    return best_num, best_score

REL_LAYOUT = "http://schemas.openxmlformats.org/officeDocument/2006/relationships/slideLayout"

def restructure_cover_slide(slide, prs, cover_image_path=None, template_prs=None):
    obj_ph = None
    for shape in slide.shapes:
        try:
            if not shape.is_placeholder:
                continue
            ph = shape.placeholder_format
            ph_type = str(ph.type)
            if ph_type not in {"CENTER_TITLE (3)", "SUBTITLE (4)", "TITLE (1)"}:
                obj_ph = shape
                break
        except Exception:
            pass
    if obj_ph is None:
        return

    txBody_src = obj_ph._element.find(f"{{{P_NS}}}txBody")
    if txBody_src is None:
        return
    src_paras = txBody_src.findall(f"{{{A_NS}}}p")
    if len(src_paras) < 1:
        return

    subtitle_para = src_paras[0]
    title_paras   = src_paras[1:] if len(src_paras) > 1 else src_paras[:1]

    title_layout = prs.slide_master.slide_layouts[0]
    for rId, rel in slide.part.rels.items():
        if REL_LAYOUT in rel.reltype:
            rel._target = title_layout.part
            break

    spTree = slide.shapes._spTree
    spTree.remove(obj_ph._element)

    next_id = max(
        (int(el.get("id", 0)) for el in spTree.iter() if el.get("id") and el.get("id").isdigit()),
        default=10
    ) + 1

    for layout_ph in title_layout.placeholders:
        ph_type = str(layout_ph.placeholder_format.type)
        if ph_type not in {"CENTER_TITLE (3)", "SUBTITLE (4)"}:
            continue

        sp_el = copy.deepcopy(layout_ph._element)

        cNvPr = sp_el.find(f".//{{{P_NS}}}cNvPr")
        if cNvPr is not None:
            cNvPr.set("id", str(next_id))
            next_id += 1

        txBody = sp_el.find(f"{{{P_NS}}}txBody")
        if txBody is None:
            continue
        for p in txBody.findall(f"{{{A_NS}}}p"):
            txBody.remove(p)

        if ph_type == "CENTER_TITLE (3)":
            for p in title_paras:
                txBody.append(copy.deepcopy(p))
        else:
            p_copy = copy.deepcopy(subtitle_para)
            for t_el in p_copy.findall(f".//{{{A_NS}}}t"):
                if t_el.text:
                    t_el.text = t_el.text.upper()
            txBody.append(p_copy)

        spTree.append(sp_el)

    if cover_image_path and os.path.exists(cover_image_path):
        pic_ph = None
        for layout in [template_prs.slide_master.slide_layouts[0] if template_prs else None,
                       prs.slide_master.slide_layouts[0]]:
            if layout is None:
                continue
            for ph in layout.placeholders:
                if str(ph.placeholder_format.type) == "PICTURE (18)":
                    pic_ph = ph
                    break
            if pic_ph:
                break

        if pic_ph is not None:
            sp_pic = pic_ph._element
            xfrm = sp_pic.find(f".//{{{A_NS}}}xfrm")
            if xfrm is not None:
                off = xfrm.find(f"{{{A_NS}}}off")
                ext = xfrm.find(f"{{{A_NS}}}ext")
                left   = int(off.get("x", 0))
                top    = int(off.get("y", 0))
                width  = int(ext.get("cx", 0))
                height = int(ext.get("cy", 0))
                slide.shapes.add_picture(cover_image_path, left, top, width, height)

def fix_cover_title(slide):
    for shape in slide.shapes:
        try:
            ph = shape.placeholder_format
        except Exception:
            continue
        if str(ph.type) != "CENTER_TITLE (3)":
            continue
        
        text = shape.text_frame.text.strip()
        text_len = len(text)
        
        target_sz = None
        if text_len > 80:
            target_sz = "3600"
        elif text_len > 60:
            target_sz = "4400"
        elif text_len > 40:
            target_sz = "4800"
            
        if target_sz:
            for rPr in shape._element.iter(f"{{{A_NS}}}rPr"):
                rPr.set("sz", target_sz)

        txBody = shape.text_frame._txBody
        bodyPr = txBody.find(f"{{{A_NS}}}bodyPr")
        if bodyPr is not None:
            for af_tag in [f"{{{A_NS}}}normAutofit", f"{{{A_NS}}}spAutoFit", f"{{{A_NS}}}noAutofit"]:
                for old in bodyPr.findall(af_tag):
                    bodyPr.remove(old)
            normAF = etree.SubElement(bodyPr, f"{{{A_NS}}}normAutofit")

_FIG_PAT = re.compile(r'insert\s+(figure|table)\s+([\d.-]+)', re.IGNORECASE)

def _add_styled_runs(p, runs, font_size_pt, font_family, font_color, template,
                     fallback_bold=False, fallback_italic=False):
    color_scheme = template.get("theme", {}).get("colorScheme", {}) if template else {}
    font_scheme = template.get("theme", {}).get("fontScheme", {}) if template else {}

    resolved_family = None
    if font_family:
        if font_family in THEME_FONTS:
            key = THEME_FONT_KEY.get(font_family, "minor")
            resolved_family = font_scheme.get(key) if font_scheme else None
        else:
            resolved_family = font_family

    resolved_color = None
    if font_color:
        resolved_color = resolve_color(font_color, color_scheme)

    for run_data in runs:
        run = p.add_run()
        run.text = run_data.get("text", "")
        run.font.size = Pt(font_size_pt)
        run.font.bold = run_data.get("bold", fallback_bold)
        run.font.italic = run_data.get("italic", fallback_italic)
        if resolved_family:
            run.font.name = resolved_family
        if resolved_color:
            try:
                run.font.color.rgb = resolved_color
            except Exception:
                pass

def insert_figure_placeholders(prs, input_dir, figures_metadata=None, template=None, include_figure_captions=True, include_table_captions=True):
    extract_dir = os.path.join(input_dir, "pdf_extracts")
    if not os.path.isdir(extract_dir):
        return []

    fig_map = {}
    for fname in os.listdir(extract_dir):
        if not fname.lower().endswith(('.png', '.jpg', '.jpeg')):
            continue
        key = os.path.splitext(fname)[0].strip().lower()
        fig_map[key] = os.path.join(extract_dir, fname)

    if not fig_map:
        return []

    caption_map = {}
    credit_map = {}
    caption_runs_map = {}
    credit_runs_map = {}
    alt_text_map = {}
    if figures_metadata:
        for fig in figures_metadata:
            name = fig.get("name")
            caption = fig.get("caption")
            credit = fig.get("credit")
            if name:
                m = re.match(r"(figure|table)\s*([\d.-]+)", name, re.IGNORECASE)
                if m:
                    key = f"{m.group(1).lower()} {m.group(2)}"
                else:
                    key = name.lower()
                if caption:
                    caption_map[key] = caption
                if credit:
                    credit_map[key] = credit
                if fig.get("captionRuns"):
                    caption_runs_map[key] = fig["captionRuns"]
                if fig.get("creditRuns"):
                    credit_runs_map[key] = fig["creditRuns"]
                if fig.get("alt_text"):
                    alt_text_map[key] = fig["alt_text"]

    used = []
    for slide_idx, slide in enumerate(prs.slides):
        replacements = []
        all_shapes = []
        def recurse(container):
            for sh in container:
                if sh.shape_type == 6:
                    try:
                        recurse(sh.shapes)
                    except Exception:
                        pass
                else:
                    all_shapes.append(sh)
        recurse(slide.shapes)

        for shape in all_shapes:
            if not shape.has_text_frame:
                continue
            text = shape.text_frame.text.strip()
            m = _FIG_PAT.search(text)
            if not m:
                continue
            fig_key = f"{m.group(1)} {m.group(2)}".lower()
            if fig_key in fig_map:
                replacements.append((shape, fig_map[fig_key], fig_key, text))
            else:
                print(f"  Slide {slide_idx + 1}: image not found for '{fig_key}' — skipping")

        for shape, img_path, fig_key, original_text in replacements:
            for sh in list(slide.shapes):
                if sh.shape_type == 13:
                    try:
                        sh._element.getparent().remove(sh._element)
                    except Exception:
                        pass

            try:
                shape_idx = list(slide.shapes).index(shape)
            except ValueError:
                shape_idx = 0

            left, top, w, h = shape.left, shape.top, shape.width, shape.height
            layout_name = None

            if template:
                try:
                    layout_name, _ = find_best_layout(slide, template)
                    for lay in template.get("slideLayouts", []):
                        if lay.get("layoutName") == layout_name:
                            for ph in lay.get("placeholders", []):
                                ph_type = ph.get("placeholder", {}).get("type")
                                if ph_type == "PICTURE (18)":
                                    pos = ph.get("position", {})
                                    sz = ph.get("size", {})
                                    if pos and sz:
                                        left = int(pos["x_pt"] * PT_TO_EMU)
                                        top = int(pos["y_pt"] * PT_TO_EMU)
                                        w = int(sz["width_pt"] * PT_TO_EMU)
                                        h = int(sz["height_pt"] * PT_TO_EMU)
                                    break
                except Exception as e:
                    print(f"  Warning: Failed to resolve template picture geometry: {e}")

            print(f"  Slide {slide_idx + 1}: placeholder dims left={left} top={top} w={w} h={h}")
            shape._element.getparent().remove(shape._element)
            
            caption_text = caption_map.get(fig_key)
            if not caption_text:
                m_orig = _FIG_PAT.search(original_text)
                if m_orig:
                    label = f"{m_orig.group(1).capitalize()} {m_orig.group(2)}"
                    rest = original_text[m_orig.end():].strip()
                    rest = re.sub(r'^[:\.\s\-\u2002\xa0]+', '', rest).strip()
                    if rest:
                        caption_text = f"{label}: {rest}"
                    else:
                        caption_text = label
                else:
                    caption_text = original_text

            credit_text = credit_map.get(fig_key)
            caption_runs = caption_runs_map.get(fig_key)
            credit_runs = credit_runs_map.get(fig_key)

            _fig_type = fig_key.split()[0] if fig_key else ""
            if _fig_type == "figure" and not include_figure_captions:
                caption_text = None
                credit_text = None
                caption_runs = None
                credit_runs = None
            elif _fig_type == "table" and not include_table_captions:
                caption_text = None
                credit_text = None
                caption_runs = None
                credit_runs = None

            font_size_pt = 18.0
            font_italic = True
            font_bold = False
            font_color = None
            font_family = "Arial"

            cred_font_size_pt = 10.0
            cred_italic = True
            cred_bold = False
            cred_font_family = "Arial"
            cred_color = None

            if caption_text and template:
                try:
                    layout_name, _ = find_best_layout(slide, template)
                    caption_props = None
                    for lay in template.get("slideLayouts", []):
                        if lay.get("layoutName") == layout_name:
                            for ph in lay.get("placeholders", []):
                                idx = ph.get("placeholder", {}).get("idx")
                                if idx in (10, 11):
                                    tb = ph.get("textBody", {})
                                    lst = tb.get("listStyle", {})
                                    lvl1 = lst.get("lvl1pPr", {})
                                    caption_props = lvl1.get("defaultRunProps", {})
                                    break
                            if caption_props:
                                break
                    if caption_props:
                        if "fontSize_pt" in caption_props:
                            font_size_pt = caption_props["fontSize_pt"]
                        if "italic" in caption_props:
                            font_italic = caption_props["italic"]
                        if "bold" in caption_props:
                            font_bold = caption_props["bold"]
                        if "fontFamily" in caption_props:
                            font_family = caption_props["fontFamily"]
                        if "color" in caption_props:
                            font_color = caption_props["color"]
                except Exception as e:
                    print(f"  Warning: Failed to resolve template caption style: {e}")

            if credit_text and template:
                try:
                    layout_name, _ = find_best_layout(slide, template)
                    credit_props = None
                    for lay in template.get("slideLayouts", []):
                        if lay.get("layoutName") == layout_name:
                            for ph in lay.get("placeholders", []):
                                idx = ph.get("placeholder", {}).get("idx")
                                if idx in (10, 11):
                                    tb = ph.get("textBody", {})
                                    lst = tb.get("listStyle", {})
                                    lvl1 = lst.get("lvl1pPr", {})
                                    credit_props = lvl1.get("defaultRunProps", {})
                                    break
                            if credit_props:
                                break
                    if credit_props:
                        if "fontSize_pt" in credit_props:
                            cred_font_size_pt = credit_props["fontSize_pt"]
                        if "italic" in credit_props:
                            cred_italic = credit_props["italic"]
                        if "bold" in credit_props:
                            cred_bold = credit_props["bold"]
                        if "fontFamily" in credit_props:
                            cred_font_family = credit_props["fontFamily"]
                        if "color" in credit_props:
                            cred_color = credit_props["color"]
                except Exception as e:
                    print(f"  Warning: Failed to resolve template credit style: {e}")

            text_width = w
            cap_left = left
            cap_width = w
            if template and layout_name:
                try:
                    for lay in template.get("slideLayouts", []):
                        if lay.get("layoutName") == layout_name:
                            for ph in lay.get("placeholders", []):
                                idx = ph.get("placeholder", {}).get("idx")
                                if idx in (10, 11) and ph.get("placeholder", {}).get("type") != "PICTURE (18)":
                                    pos = ph.get("position", {})
                                    sz = ph.get("size", {})
                                    if pos and sz:
                                        cap_left = int(pos["x_pt"] * PT_TO_EMU)
                                        cap_width = int(sz["width_pt"] * PT_TO_EMU)
                                        text_width = cap_width
                                        break
                except Exception:
                    pass

            reserved_h = Pt(0)
            cap_height = Pt(0)
            cred_height = Pt(0)

            if caption_text:
                cap_chars_per_line = max(10, int((text_width / PT_TO_EMU) / (font_size_pt * 0.38)))
                cap_lines = (len(caption_text) + cap_chars_per_line - 1) // cap_chars_per_line
                cap_height = Pt(cap_lines * (font_size_pt + 4) + 12)
                reserved_h += cap_height + Pt(10)

            if credit_text:
                cred_chars_per_line = max(10, int((text_width / PT_TO_EMU) / (cred_font_size_pt * 0.38)))
                cred_lines = (len(credit_text) + cred_chars_per_line - 1) // cred_chars_per_line
                cred_height = Pt(cred_lines * (cred_font_size_pt + 4) + 12)
                reserved_h += cred_height + Pt(5)

            max_pic_h = h
            if reserved_h > 0:
                max_pic_h = max(Pt(100), h - reserved_h)

            pic = None
            if w <= 0 or max_pic_h <= 0:
                pic = slide.shapes.add_picture(img_path, left, top)
                max_w = int(prs.slide_width * 0.60)
                if pic.width > max_w:
                    ratio = max_w / pic.width
                    pic.width  = max_w
                    pic.height = int(pic.height * ratio)
            else:
                pic = slide.shapes.add_picture(img_path, left, top)
                orig_w, orig_h = pic.width, pic.height
                ratio = min(w / orig_w, max_pic_h / orig_h)
                pic.width = int(orig_w * ratio)
                pic.height = int(orig_h * ratio)
                pic.left = left + int((w - pic.width) / 2)
                pic.top = top + int((max_pic_h - pic.height) / 2)

            if pic:
                try:
                    cNvPr = pic._element[0][0]
                    if fig_key in alt_text_map:
                        cNvPr.set('descr', alt_text_map[fig_key])
                    else:
                        cNvPr.attrib.pop('descr', None)
                except Exception:
                    pass

            if caption_text and pic:
                cap_left = pic.left
                cap_width = (left + w) - pic.left

                adjusted_font_size = font_size_pt
                adjusted_cred_font_size = cred_font_size_pt
                
                cap_top = pic.top + pic.height + Pt(10)
                
                total_chars = len(caption_text or '') + len(credit_text or '')
                if total_chars > 0:
                    box_width_in = float(cap_width) / 914400.0
                    max_total_height_in = (float(prs.slide_height) - float(cap_top)) / 914400.0 - 0.1
                    if max_total_height_in <= 0.5:
                        max_total_height_in = 1.8
                        
                    avg_font_size = (adjusted_font_size + adjusted_cred_font_size) / 2.0
                    line_height_in = (avg_font_size / 72.0) * 1.25
                    char_width_in = (avg_font_size / 72.0) * 0.43
                    
                    chars_per_line = max(1, int(box_width_in / char_width_in))
                    max_lines = max(1, int(max_total_height_in / line_height_in))
                    total_capacity = chars_per_line * max_lines
                    
                    if total_chars > total_capacity:
                        scale_factor = (float(total_capacity) / total_chars) ** 0.5
                        adjusted_font_size = max(7.0, adjusted_font_size * scale_factor)
                        adjusted_cred_font_size = max(6.0, adjusted_cred_font_size * scale_factor)

                cap_chars_per_line = max(10, int((cap_width / PT_TO_EMU) / (adjusted_font_size * 0.38)))
                cap_lines = (len(caption_text) + cap_chars_per_line - 1) // cap_chars_per_line
                cap_height = Pt(cap_lines * (adjusted_font_size + 4) + 12)

                if credit_text:
                    cred_chars_per_line = max(10, int((cap_width / PT_TO_EMU) / (adjusted_cred_font_size * 0.38)))
                    cred_lines = (len(credit_text) + cred_chars_per_line - 1) // cred_chars_per_line
                    cred_height = Pt(cred_lines * (adjusted_cred_font_size + 4) + 12)
                else:
                    cred_height = Pt(0)

                total_text_height = cap_height + (cred_height + Pt(5) if credit_text else Pt(0))
                if cap_top + total_text_height > prs.slide_height:
                    cap_top = prs.slide_height - total_text_height - Pt(5)
                    if cap_top < pic.top + pic.height:
                        cap_top = pic.top + pic.height + Pt(2)

                txBox = slide.shapes.add_textbox(cap_left, cap_top, cap_width, cap_height)
                tf = txBox.text_frame
                tf.word_wrap = True
                try:
                    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                except Exception:
                    pass
                p = tf.paragraphs[0]
                if caption_runs:
                    _add_styled_runs(p, caption_runs, adjusted_font_size,
                                     font_family, font_color, template,
                                     fallback_bold=font_bold, fallback_italic=font_italic)
                else:
                    run = p.add_run()
                    run.text = caption_text
                    run.font.size = Pt(adjusted_font_size)
                    run.font.bold = font_bold
                    run.font.italic = font_italic
                    if font_family:
                        if font_family in THEME_FONTS:
                            key = THEME_FONT_KEY.get(font_family, "minor")
                            font_scheme = template.get("theme", {}).get("fontScheme", {})
                            resolved_fn = font_scheme.get(key) if font_scheme else None
                            if resolved_fn:
                                run.font.name = resolved_fn
                        else:
                            run.font.name = font_family
                    if font_color:
                        color_scheme = template.get("theme", {}).get("colorScheme", {})
                        rgb = resolve_color(font_color, color_scheme)
                        if rgb:
                            try:
                                run.font.color.rgb = rgb
                            except Exception:
                                pass

                if credit_text:
                    cred_left = cap_left
                    cred_width = cap_width
                    cred_top = cap_top + cap_height + Pt(5)

                    txBoxCred = slide.shapes.add_textbox(cred_left, cred_top, cred_width, cred_height)
                    tf_cred = txBoxCred.text_frame
                    tf_cred.word_wrap = True
                    try:
                        tf_cred.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    except Exception:
                        pass
                    p_cred = tf_cred.paragraphs[0]
                    if credit_runs:
                        _add_styled_runs(p_cred, credit_runs, adjusted_cred_font_size,
                                         cred_font_family, cred_color, template,
                                         fallback_bold=cred_bold, fallback_italic=cred_italic)
                    else:
                        run_cred = p_cred.add_run()
                        run_cred.text = credit_text
                        run_cred.font.size = Pt(adjusted_cred_font_size)
                        run_cred.font.italic = cred_italic
                        run_cred.font.bold = cred_bold
                        if cred_font_family in THEME_FONTS:
                            key = THEME_FONT_KEY.get(cred_font_family, "minor")
                            font_scheme = template.get("theme", {}).get("fontScheme", {})
                            resolved_fn = font_scheme.get(key) if font_scheme else None
                            if resolved_fn:
                                run_cred.font.name = resolved_fn
                        else:
                            run_cred.font.name = cred_font_family

                    if cred_color:
                        color_scheme = template.get("theme", {}).get("colorScheme", {})
                        rgb = resolve_color(cred_color, color_scheme)
                        if rgb:
                            try:
                                p_cred.font.color.rgb = rgb
                            except Exception:
                                pass

            fname = os.path.basename(img_path)
            used.append({
                "dest_name": fname,
                "slideIndex": slide_idx,
                "shapeIndex": shape_idx
            })
            print(f"  Slide {slide_idx + 1}: replaced placeholder → {fname}")

    return used

def insert_extracted_images(prs, input_dir, skip=None):
    manifest_path = os.path.join(input_dir, "pdf_extracts", "manifest.json")
    if not os.path.exists(manifest_path):
        return
    with open(manifest_path) as f:
        manifest = json.load(f)
    extracts = [e for e in manifest.get("extracts", []) if e.endswith(".png")]
    if not extracts:
        return

    extract_dir  = os.path.join(input_dir, "pdf_extracts")
    blank_layout = prs.slide_master.slide_layouts[6]
    SLIDE_W      = prs.slide_width
    SLIDE_H      = prs.slide_height

    inserted = 0
    for fname in extracts:
        if skip and fname in skip:
            continue
        fpath = os.path.join(extract_dir, fname)
        if not os.path.exists(fpath):
            print(f"  [warn] extracted image not found, skipping: {fname}")
            continue
        slide = prs.slides.add_slide(blank_layout)
        pic   = slide.shapes.add_picture(fpath, 0, 0)
        scale = min(SLIDE_W / pic.width, SLIDE_H / pic.height)
        pic.width  = int(pic.width  * scale)
        pic.height = int(pic.height * scale)
        pic.left   = (SLIDE_W - pic.width)  // 2
        pic.top    = (SLIDE_H - pic.height) // 2
        inserted  += 1

    if inserted:
        print(f"  Appended {inserted} extracted image slide(s) from pdf_extracts/")

def find_cover_image(input_path, explicit_cover=None):
    if explicit_cover and os.path.exists(explicit_cover):
        return explicit_cover

    input_dir  = os.path.dirname(os.path.abspath(input_path))
    input_base = os.path.basename(input_path)
    isbn = input_base.split("_")[0] if "_" in input_base else ""

    img_exts = (".jpg", ".jpeg", ".png", ".bmp", ".gif")
    candidates = []
    for fname in os.listdir(input_dir):
        if not fname.lower().endswith(img_exts):
            continue
        fpath = os.path.join(input_dir, fname)
        if isbn and fname.startswith(isbn):
            candidates.insert(0, fpath)
        elif fname.lower().startswith("cover"):
            candidates.append(fpath)

    return candidates[0] if candidates else None

def fix_overflowing_textboxes(prs):
    import math
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame or not shape.text_frame.text.strip():
                continue
            
            try:
                if shape.is_placeholder and str(shape.placeholder_format.type) in ["TITLE (1)", "CENTER_TITLE (3)"]:
                    continue
            except Exception:
                pass
            
            tf = shape.text_frame
            width_pt = shape.width / 12700
            height_pt = shape.height / 12700
            
            txBody = tf._txBody
            bodyPr = txBody.find(f"{{{A_NS}}}bodyPr")
            lIns = 10.0
            rIns = 10.0
            tIns = 5.0
            bIns = 5.0
            if bodyPr is not None:
                if bodyPr.get("lIns") is not None: lIns = int(bodyPr.get("lIns")) / 12700
                if bodyPr.get("rIns") is not None: rIns = int(bodyPr.get("rIns")) / 12700
                if bodyPr.get("tIns") is not None: tIns = int(bodyPr.get("tIns")) / 12700
                if bodyPr.get("bIns") is not None: bIns = int(bodyPr.get("bIns")) / 12700
            
            avail_w = max(50.0, width_pt - lIns - rIns)
            avail_h = max(20.0, height_pt - tIns - bIns)
            
            total_est_h = 0.0
            para_details = []
            
            for para in tf.paragraphs:
                pt_sz = 20.0
                if para.runs:
                    for r in para.runs:
                        if r.font.size is not None:
                            pt_sz = r.font.size.pt
                            break
                elif para.font.size is not None:
                    pt_sz = para.font.size.pt

                sb = 0.0
                sa = Pt(6).pt
                if para.space_before is not None:
                    sb = para.space_before.pt
                if para.space_after is not None:
                    sa = para.space_after.pt
                ls = 1.15
                if para.line_spacing is not None:
                    if isinstance(para.line_spacing, float):
                        ls = para.line_spacing
                    else:
                        ls = para.line_spacing.pt / pt_sz

                text = para.text
                if not text.strip():
                    total_est_h += (pt_sz * ls) + sb + sa
                    continue

                char_w = pt_sz * 0.38
                est_text_w = len(text) * char_w
                lines = max(1.0, math.ceil(est_text_w / avail_w))

                para_h = lines * (pt_sz * ls) + sb + sa
                total_est_h += para_h
                para_details.append((para, pt_sz, para_h))
            
            if total_est_h > avail_h and para_details:
                scale = avail_h / total_est_h
                for para, orig_sz, _ in para_details:
                    scaled = orig_sz * scale
                    new_sz = max(10.0, math.floor(scaled / 2) * 2)
                    for r in para.runs:
                        r.font.size = Pt(new_sz)
                    if not para.runs:
                        para.font.size = Pt(new_sz)

def convert(input_path, template_style_path, output_path, apply_geometry=True, cover_image_path=None, figures_metadata=None, include_figure_captions=True, include_table_captions=True):
    if not os.path.exists(input_path):
        print(f"Error: input file '{input_path}' not found.")
        sys.exit(1)
    if not os.path.exists(template_style_path):
        print(f"Error: template style file '{template_style_path}' not found.")
        sys.exit(1)

    cover_image_path = find_cover_image(input_path, cover_image_path)
    if cover_image_path:
        print(f"  Cover image: {os.path.basename(cover_image_path)}")

    with open(template_style_path) as f:
        template = json.load(f)

    color_scheme    = template["theme"]["colorScheme"]
    font_scheme     = template["theme"].get("fontScheme", {})
    master_styles   = template["masterTextStyles"]
    master_bg_fill  = template.get("masterBackground", {}).get("fill")

    template_pptx_path = template.get("_meta", {}).get("source")
    template_prs = None
    if template_pptx_path and os.path.exists(template_pptx_path):
        template_prs = Presentation(template_pptx_path)

    prs = Presentation(input_path)

    strip_copyright_from_masters(prs)

    tmpl_w = int(template["presentation"]["width_pt"]  * PT_TO_EMU)
    tmpl_h = int(template["presentation"]["height_pt"] * PT_TO_EMU)
    prs.slide_width  = tmpl_w
    prs.slide_height = tmpl_h

    for slide_idx, slide in enumerate(prs.slides):
        if slide_idx == 0:
            restructure_cover_slide(slide, prs,
                                    cover_image_path=cover_image_path,
                                    template_prs=template_prs)

        layout_name, layout_score = find_best_layout(slide, template)
        template_slide_num, slide_score = find_best_template_slide(slide, template)

        print(f"  Slide {slide_idx + 1} → layout: {layout_name} (score {layout_score}), "
              f"template slide: {template_slide_num} (score {slide_score})")

        real_ph_indices = {
            s.placeholder_format.idx for s in slide.shapes
            if s.is_placeholder and s.placeholder_format is not None
        }
        nonph_text_shapes = [
            s for s in slide.shapes
            if not s.is_placeholder and s.has_text_frame
            and any(p.text.strip() for p in s.text_frame.paragraphs)
            and not _FIG_PAT.search(s.text_frame.text)
        ]

        layout_bg_fill = None
        for layout in template["slideLayouts"]:
            if layout["layoutName"] == layout_name:
                layout_bg_fill = layout.get("background", {}).get("fill")
                break
        apply_slide_background(slide, layout_bg_fill or master_bg_fill, color_scheme)

        add_decorative_shapes(slide, layout_name, template, color_scheme, template_prs)

        present_indices = {sh.placeholder_format.idx for sh in slide.shapes if sh.is_placeholder and sh.placeholder_format is not None}

        _layout_for_slide = next(
            (l for l in template["slideLayouts"] if l["layoutName"] == layout_name), {}
        )

        for shape in slide.shapes:
            if not shape.is_placeholder or not shape.has_text_frame:
                continue

            ph_idx  = shape.placeholder_format.idx
            ph_type = str(shape.placeholder_format.type)
            is_title = any(t in ph_type for t in TITLE_PH_TYPES)
            is_title_or_subtitle = is_title or "SUBTITLE" in ph_type

            if layout_name == "slideLayout5" and ph_idx == 1 and 2 not in present_indices:
                ph_idx = 2

            layout_ph   = get_layout_ph(template, layout_name, ph_idx)
            master_sh   = get_master_shape(template, ph_idx)
            slide_style, template_sh = get_template_slide_ph(template, template_slide_num, ph_idx)

            if apply_geometry:
                def _first_geo():
                    for src in (template_sh, layout_ph, master_sh):
                        if src and src.get("position") and src.get("size"):
                            return src
                    return None
                geo = _first_geo()
                orig_ph_idx = shape.placeholder_format.idx
                other_ph_indices = present_indices - {orig_ph_idx}
                if geo and geo.get("position") and geo.get("size"):
                    body_y = geo["position"]["y_pt"]
                    body_h = geo["size"]["height_pt"]
                    for _adj in _layout_for_slide.get("placeholders", []):
                        _adj_idx = _adj.get("placeholder", {}).get("idx")
                        if _adj_idx in other_ph_indices or _adj_idx == ph_idx:
                            continue
                        _adj_y = (_adj.get("position") or {}).get("y_pt")
                        _adj_h = (_adj.get("size") or {}).get("height_pt")
                        if _adj_y is None or _adj_h is None or _adj_y >= body_y:
                            continue
                        if abs((_adj_y + _adj_h) - body_y) <= 15:
                            geo = {
                                "position": {"x_pt": geo["position"]["x_pt"], "y_pt": _adj_y},
                                "size": {"width_pt": geo["size"]["width_pt"],
                                         "height_pt": (body_y + body_h) - _adj_y}
                            }
                            break
                apply_shape_geometry(shape, geo)

            def _first_with(key):
                for src in (template_sh, layout_ph, master_sh):
                    if src and key in src:
                        return src
                return None

            apply_shape_fill(shape, _first_with("fill"), color_scheme)
            apply_shape_border(shape, _first_with("border"), color_scheme)

            bp = (layout_ph or {}).get("textBody", {}).get("bodyProperties", {})
            if not bp and master_sh:
                bp = master_sh.get("textBody", {}).get("bodyProperties", {})
            apply_body_properties(shape.text_frame, bp)

            for para in shape.text_frame.paragraphs:
                level = para.level or 0

                if is_title:
                    master_style = get_master_style_for_title(master_styles)
                else:
                    master_style = get_master_style_for_body(master_styles, level)

                layout_style = get_layout_style(template, layout_name, ph_idx, level)

                final_style = merge(master_style, layout_style, slide_style)

                if is_title_or_subtitle:
                    for _bk in ('bulletChar', 'bulletFont', 'bulletColor',
                                'marginLeft_pt', 'indent_pt'):
                        final_style.pop(_bk, None)

                apply_para_style(para, final_style, color_scheme)

                if is_title:
                    pattern = re.compile(r"(\(\d+\s+[oO][fF]\s+\d+\)|\(\s*[cC][oO][nN][tT][a-zA-Z.]*\s*\))")
                    text = para.text
                    parts = pattern.split(text)
                    if len(parts) > 1:
                        para.text = ""
                        for part in parts:
                            if not part:
                                continue
                            run = para.add_run()
                            run.text = part
                            apply_run_style(run, final_style, color_scheme, font_scheme)
                            if pattern.match(part):
                                run.font.size = Pt(20)
                    else:
                        for run in para.runs:
                            apply_run_style(run, final_style, color_scheme, font_scheme)
                else:
                    for run in para.runs:
                        apply_run_style(run, final_style, color_scheme, font_scheme)

        _MIN_BODY_SLOT_PT = 60
        _body_slot_pool = [
            ph.get("placeholder", {}).get("idx")
            for ph in sorted(
                _layout_for_slide.get("placeholders", []),
                key=lambda p: p.get("placeholder", {}).get("idx", 999)
            )
            if ph.get("placeholder", {}).get("idx") not in real_ph_indices
            and (ph.get("size") or {}).get("height_pt", 999) >= _MIN_BODY_SLOT_PT
        ]

        virtual_start = (max(real_ph_indices) + 1) if real_ph_indices else 2
        for v_offset, shape in enumerate(nonph_text_shapes):
            if v_offset < len(_body_slot_pool):
                virtual_idx = _body_slot_pool[v_offset]
            else:
                virtual_idx = virtual_start + v_offset
            layout_ph_v = get_layout_ph(template, layout_name, virtual_idx)
            master_sh_v = get_master_shape(template, virtual_idx)

            if not layout_ph_v and not master_sh_v:
                continue

            if apply_geometry:
                def _first_geo_v():
                    for src in (layout_ph_v, master_sh_v):
                        if src and src.get("position") and src.get("size"):
                            return src
                    return None
                geo = _first_geo_v()
                if geo and geo.get("position") and geo.get("size"):
                    body_y = geo["position"]["y_pt"]
                    body_h = geo["size"]["height_pt"]
                    for _adj in _layout_for_slide.get("placeholders", []):
                        _adj_idx = _adj.get("placeholder", {}).get("idx")
                        if _adj_idx in real_ph_indices or _adj_idx == virtual_idx:
                            continue
                        _adj_y = (_adj.get("position") or {}).get("y_pt")
                        _adj_h = (_adj.get("size") or {}).get("height_pt")
                        if _adj_y is None or _adj_h is None or _adj_y >= body_y:
                            continue
                        if abs((_adj_y + _adj_h) - body_y) <= 15:
                            merged_y = _adj_y
                            geo = {
                                "position": {"x_pt": geo["position"]["x_pt"], "y_pt": merged_y},
                                "size": {"width_pt": geo["size"]["width_pt"],
                                         "height_pt": (body_y + body_h) - merged_y}
                            }
                            break
                apply_shape_geometry(shape, geo)

            bp_v = (layout_ph_v or {}).get("textBody", {}).get("bodyProperties", {})
            if not bp_v:
                master_body = get_master_shape(template, 1)
                if master_body:
                    bp_v = master_body.get("textBody", {}).get("bodyProperties", {})
            apply_body_properties(shape.text_frame, bp_v)

            for para in shape.text_frame.paragraphs:
                level = para.level or 0
                master_style = get_master_style_for_body(master_styles, level)
                layout_style = get_layout_style(template, layout_name, virtual_idx, level)
                final_style = merge(master_style, layout_style, {})

                apply_para_style(para, final_style, color_scheme)
                for run in para.runs:
                    apply_run_style(run, final_style, color_scheme, font_scheme)

        if slide_idx == 0:
            fix_cover_title(slide)

    input_dir = os.path.dirname(os.path.abspath(input_path))
    used_figs = insert_figure_placeholders(prs, input_dir, figures_metadata, template=template, include_figure_captions=include_figure_captions, include_table_captions=include_table_captions)

    skip_names = {x["dest_name"] for x in used_figs}
    insert_extracted_images(prs, input_dir, skip=skip_names)

    fix_overflowing_textboxes(prs)

    prs.save(output_path)
    return used_figs
