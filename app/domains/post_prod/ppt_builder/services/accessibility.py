from pptx import Presentation

MIN_FONT_SIZE = 18
GENERIC_LINK_TEXT = ["click here", "here", "link"]

def get_relative_luminance(color_rgb):
    r, g, b = color_rgb[0] / 255.0, color_rgb[1] / 255.0, color_rgb[2] / 255.0
    def adjust(val):
        return val / 12.92 if val <= 0.03928 else ((val + 0.055) / 1.055) ** 2.4
    return 0.2126 * adjust(r) + 0.7152 * adjust(g) + 0.0722 * adjust(b)

def get_contrast_ratio(rgb1, rgb2):
    l1 = get_relative_luminance(rgb1)
    l2 = get_relative_luminance(rgb2)
    lighter = max(l1, l2)
    darker = min(l1, l2)
    return (lighter + 0.05) / (darker + 0.05)

def check_ppt_accessibility(file_path):
    prs = Presentation(file_path)
    slide_w_pt = prs.slide_width / 12700
    slide_h_pt = prs.slide_height / 12700
    issues = []

    for slide_index, slide in enumerate(prs.slides, start=1):
        if not slide.shapes.title or not slide.shapes.title.text.strip():
            issues.append({"slide": slide_index, "category": "Missing Slide Title", "severity": "Error", "detail": "Slide is missing a title shape."})

        for shape in slide.shapes:
            if hasattr(shape, "text"):
                if not shape.text.strip():
                    issues.append({"slide": slide_index, "category": "Empty Text Box", "severity": "Tip", "detail": "An empty text box was found on the slide."})

            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_size = run.font.size.pt
                            if font_size < MIN_FONT_SIZE:
                                issues.append({"slide": slide_index, "category": "Small Font Size", "severity": "Tip", "detail": f"Small font ({font_size}pt) used in text '{run.text}'"})

                        if run.hyperlink.address:
                            if run.text.lower() in GENERIC_LINK_TEXT:
                                issues.append({"slide": slide_index, "category": "Generic Hyperlink Text", "severity": "Warning", "detail": f"Generic hyperlink text '{run.text}' used."})

            if slide.shapes.title and len(slide.shapes) > 0:
                if slide.shapes[0] != slide.shapes.title:
                    issues.append({"slide": slide_index, "category": "Reading Order Issue", "severity": "Warning", "detail": "Title is not the first element read by screen readers."})

            if shape.shape_type == 13:  # Picture
                try:
                    descr = (shape._element.nvPicPr.cNvPr.get('descr') or '').strip()
                except Exception:
                    descr = ''
                if not descr:
                    issues.append({"slide": slide_index, "category": "Missing Alt Text", "severity": "Error", "detail": "Image is missing meaningful alt text."})

            if shape.has_table:
                table = shape.table
                first_row_empty = all(
                    cell.text.strip() == "" for cell in table.rows[0].cells
                )
                header_style_disabled = not getattr(table, "first_row", True)
                if first_row_empty:
                    issues.append({"slide": slide_index, "category": "Missing Table Header", "severity": "Error", "detail": "Table header row is missing or empty."})
                elif header_style_disabled:
                    issues.append({"slide": slide_index, "category": "Missing Table Header", "severity": "Warning", "detail": "Table has header text, but header row styling is disabled."})
                    
            if hasattr(shape, "fill") and shape.fill.type == 1: # SOLID
                try:
                    bg_color = shape.fill.fore_color.rgb
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                if run.font.color and run.font.color.type == 1: # RGB
                                    fg_color = run.font.color.rgb
                                    font_size = run.font.size.pt if run.font.size else 14
                                    is_large_text = font_size >= 18 or (font_size >= 14 and run.font.bold)
                                    ratio = get_contrast_ratio(bg_color, fg_color)
                                    required = 3.0 if is_large_text else 4.5
                                    if ratio < required:
                                        issues.append({
                                            "slide": slide_index,
                                            "category": "Poor Color Contrast",
                                            "severity": "Error",
                                            "detail": f"Contrast ratio is {ratio:.1f}:1, which is below the WCAG AA requirement of {required}:1 for text '{run.text}'."
                                        })
                except Exception:
                    pass

            is_media = (
                (hasattr(shape, "shape_type") and int(shape.shape_type) in [16, 24]) or 
                "video" in shape.name.lower() or 
                "audio" in shape.name.lower() or 
                "media" in shape.name.lower()
            )
            if is_media:
                issues.append({
                    "slide": slide_index,
                    "category": "Missing Media Subtitles",
                    "severity": "Warning",
                    "detail": f"Media element '{shape.name}' was found. Ensure it has subtitles, closed captions, or a text transcript."
                })

            if shape.has_text_frame and shape.text_frame.text.strip():
                import math
                tf = shape.text_frame
                width_pt = shape.width / 12700
                height_pt = shape.height / 12700

                lIns = 91440 / 12700
                rIns = 91440 / 12700
                tIns = 45720 / 12700
                bIns = 45720 / 12700
                bodyPr = tf._txBody.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
                sp_autofit = False
                norm_autofit = False
                if bodyPr is not None:
                    if bodyPr.get("lIns") is not None: lIns = int(bodyPr.get("lIns")) / 12700
                    if bodyPr.get("rIns") is not None: rIns = int(bodyPr.get("rIns")) / 12700
                    if bodyPr.get("tIns") is not None: tIns = int(bodyPr.get("tIns")) / 12700
                    if bodyPr.get("bIns") is not None: bIns = int(bodyPr.get("bIns")) / 12700
                    sp_autofit   = bodyPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}spAutoFit") is not None
                    norm_autofit = bodyPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}normAutofit") is not None

                if sp_autofit:
                    pass
                else:
                    base_avail_w = max(10.0, width_pt - lIns - rIns)
                    avail_h = max(1.0, height_pt - tIns - bIns)

                    total_est_h = 0.0
                    for para in tf.paragraphs:
                        pt_sz = 18.0
                        if para.runs:
                            for r in para.runs:
                                if r.font.size is not None:
                                    pt_sz = r.font.size.pt
                                    break
                        elif para.font.size is not None:
                            pt_sz = para.font.size.pt

                        text = para.text
                        if not text.strip():
                            continue

                        pPr = para._p.find("{http://schemas.openxmlformats.org/drawingml/2006/main}pPr")
                        mar_l_pt = 0.0
                        if pPr is not None and pPr.get("marL"):
                            mar_l_pt = int(pPr.get("marL")) / 12700
                        para_avail_w = max(10.0, base_avail_w - mar_l_pt)

                        char_w = pt_sz * 0.38
                        lines = max(1.0, math.ceil(len(text) * char_w / para_avail_w))

                        sb = para.space_before.pt if para.space_before else 0.0
                        sa = para.space_after.pt if para.space_after else 6.0
                        ls = 1.15
                        if para.line_spacing is not None:
                            ls = para.line_spacing if isinstance(para.line_spacing, float) else para.line_spacing.pt / pt_sz

                        total_est_h += lines * (pt_sz * ls) + sb + sa

                    threshold = 0.0 if norm_autofit else 5.0
                    if total_est_h > avail_h + threshold:
                        if norm_autofit:
                            issues.append({
                                "slide": slide_index,
                                "category": "Text Shrinks to Fit",
                                "severity": "Warning",
                                "detail": f"Text in '{shape.name}' is auto-shrunk to fit (requires ~{total_est_h:.1f}pt, shape height {height_pt:.1f}pt). Font may become unreadably small."
                            })
                        else:
                            issues.append({
                                "slide": slide_index,
                                "category": "Text Overflow",
                                "severity": "Warning",
                                "detail": f"Text in '{shape.name}' potentially overflows shape bounds (requires ~{total_est_h:.1f}pt, shape height {height_pt:.1f}pt)."
                            })

        for shape in slide.shapes:
            if not shape.width or not shape.height:
                continue
            has_content = (shape.has_text_frame and shape.text_frame.text.strip()) or shape.shape_type == 13
            if not has_content:
                continue
            if shape.has_text_frame:
                bodyPr = shape.text_frame._txBody.find("{http://schemas.openxmlformats.org/drawingml/2006/main}bodyPr")
                if bodyPr is not None and bodyPr.find("{http://schemas.openxmlformats.org/drawingml/2006/main}spAutoFit") is not None:
                    continue
            left_pt   = shape.left / 12700
            top_pt    = shape.top / 12700
            right_pt  = left_pt + shape.width / 12700
            bottom_pt = top_pt + shape.height / 12700
            if right_pt > slide_w_pt + 1:
                issues.append({
                    "slide": slide_index,
                    "category": "Shape Off Slide",
                    "severity": "Warning",
                    "detail": f"Shape '{shape.name}' extends {right_pt - slide_w_pt:.1f}pt beyond the right edge of the slide — content will be clipped."
                })
            if bottom_pt > slide_h_pt + 1:
                issues.append({
                    "slide": slide_index,
                    "category": "Shape Off Slide",
                    "severity": "Warning",
                    "detail": f"Shape '{shape.name}' extends {bottom_pt - slide_h_pt:.1f}pt below the bottom edge of the slide — content will be clipped."
                })

        spatial_shapes = []
        for shape in slide.shapes:
            if not shape.width or not shape.height:
                continue
            w_pt = shape.width / 12700
            h_pt = shape.height / 12700
            if w_pt > 864 or h_pt > 486:
                continue
            has_text = shape.has_text_frame and shape.text_frame.text.strip()
            is_pic = (shape.shape_type == 13)
            if has_text or is_pic:
                l_pt = shape.left / 12700
                t_pt = shape.top / 12700
                spatial_shapes.append({
                    "name": shape.name,
                    "is_pic": is_pic,
                    "box": (l_pt, t_pt, l_pt + w_pt, t_pt + h_pt)
                })

        for idx1 in range(len(spatial_shapes)):
            for idx2 in range(idx1 + 1, len(spatial_shapes)):
                s1 = spatial_shapes[idx1]
                s2 = spatial_shapes[idx2]
                if not (s1["is_pic"] ^ s2["is_pic"]):
                    continue
                box1 = s1["box"]
                box2 = s2["box"]

                int_l = max(box1[0], box2[0])
                int_t = max(box1[1], box2[1])
                int_r = min(box1[2], box2[2])
                int_b = min(box1[3], box2[3])

                if int_r > int_l and int_b > int_t:
                    overlap_w = int_r - int_l
                    overlap_h = int_b - int_t
                    if overlap_w > 15.0 and overlap_h > 15.0:
                        issues.append({
                            "slide": slide_index,
                            "category": "Shape Overlap",
                            "severity": "Warning",
                            "detail": f"Shape '{s1['name']}' overlaps with shape '{s2['name']}' by {overlap_w:.1f}pt x {overlap_h:.1f}pt."
                        })

    title_map = {}
    for slide_index, slide in enumerate(prs.slides, start=1):
        if slide.shapes.title and slide.shapes.title.text.strip():
            title_text = slide.shapes.title.text.strip().lower()
            if title_text not in title_map:
                title_map[title_text] = []
            title_map[title_text].append(slide_index)

    for title_text, slide_indices in title_map.items():
        if len(slide_indices) > 1:
            for idx in slide_indices:
                other_slides = [str(i) for i in slide_indices if i != idx]
                issues.append({
                    "slide": idx,
                    "category": "Duplicate Slide Title",
                    "severity": "Warning",
                    "detail": f"This slide has the same title as slide(s): {', '.join(other_slides)}."
                })

    slide_contents = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                texts.append(shape.text_frame.text.strip())
            elif shape.shape_type == 13:
                texts.append(f"[Image:{shape.name}]")
        content_sig = "||".join(sorted(texts))
        slide_contents.append((slide_index, content_sig))

    content_map = {}
    for idx, sig in slide_contents:
        if sig:
            if sig not in content_map:
                content_map[sig] = []
            content_map[sig].append(idx)

    for sig, slide_indices in content_map.items():
        if len(slide_indices) > 1:
            for idx in slide_indices:
                other_slides = [str(i) for i in slide_indices if i != idx]
                issues.append({
                    "slide": idx,
                    "category": "Duplicate Slide Content",
                    "severity": "Warning",
                    "detail": f"This slide has identical content to slide(s): {', '.join(other_slides)}."
                })

    return issues
