from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from lxml import etree
import json

INVALID_PH_IDX = 4294967295
_extraction_warnings = []

NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}

EMU_PER_PT = 12700

def emu_to_pt(emu):
    if emu is None:
        return None
    return round(emu / EMU_PER_PT, 2)

def hundredths_to_pt(val):
    if val is None:
        return None
    return round(val / 100, 1)

def thousandths_to_pct(val):
    if val is None:
        return None
    return round(val / 1000, 1)

def parse_color_element(elem):
    if elem is None:
        return None
    srgb = elem.find(".//a:srgbClr", NS)
    if srgb is not None:
        return "#" + srgb.get("val", "")
    scheme = elem.find(".//a:schemeClr", NS)
    if scheme is not None:
        return "scheme:" + scheme.get("val", "")
    sys_clr = elem.find(".//a:sysClr", NS)
    if sys_clr is not None:
        return "#" + (sys_clr.get("lastClr") or sys_clr.get("val", ""))
    return None

def color_from_font(font):
    try:
        if font.color and font.color.type is not None:
            rgb = font.color.rgb
            return "#" + str(rgb)
    except Exception:
        pass
    rPr = font._element
    solid = rPr.find(".//a:solidFill", NS) if rPr is not None else None
    return parse_color_element(solid)

def extract_run_props(run):
    font = run.font
    rPr = font._element
    props = {}

    sz = rPr.get("sz") if rPr is not None else None
    props["fontSize_pt"] = hundredths_to_pt(int(sz)) if sz else (font.size.pt if font.size else None)

    props["fontFamily"] = font.name
    props["bold"] = font.bold
    props["italic"] = font.italic
    props["underline"] = font.underline if font.underline is not True else True

    color = color_from_font(font)
    if color:
        props["color"] = color

    if rPr is not None:
        kern = rPr.get("kern")
        if kern:
            props["kerning_pt"] = hundredths_to_pt(int(kern))
        spc = rPr.get("spc")
        if spc:
            props["charSpacing_pt"] = hundredths_to_pt(int(spc))
        lang = rPr.get("lang")
        if lang:
            props["language"] = lang

    return {k: v for k, v in props.items() if v is not None}

def extract_para_props(para):
    pPr = para._p.find("a:pPr", NS)
    props = {}
    props["level"] = getattr(para, "level", 0)

    align_map = {
        PP_ALIGN.LEFT: "l", PP_ALIGN.CENTER: "ctr", PP_ALIGN.RIGHT: "r",
        PP_ALIGN.JUSTIFY: "just", PP_ALIGN.DISTRIBUTE: "dist",
    }
    if para.alignment is not None:
        props["alignment"] = align_map.get(para.alignment, str(para.alignment))

    if pPr is not None:
        indent = pPr.get("indent")
        if indent:
            props["indent_pt"] = emu_to_pt(int(indent))
        marL = pPr.get("marL")
        if marL:
            props["marginLeft_pt"] = emu_to_pt(int(marL))

        for tag, key in [("a:spcBef", "spaceBefore"), ("a:spcAft", "spaceAfter"), ("a:lnSpc", "lineSpacing")]:
            el = pPr.find(tag, NS)
            if el is not None:
                pts = el.find("a:spcPts", NS)
                pct = el.find("a:spcPct", NS)
                if pts is not None:
                    props[key + "_pt"] = hundredths_to_pt(int(pts.get("val", 0)))
                elif pct is not None:
                    props[key + "_pct"] = thousandths_to_pct(int(pct.get("val", 0)))

    return {k: v for k, v in props.items() if v is not None}

def extract_text_body(tf):
    result = {}
    txBody = tf._txBody

    bodyPr = txBody.find("a:bodyPr", NS)
    if bodyPr is not None:
        bp = {}
        for attr, key in [("vert", "textDirection"), ("anchor", "verticalAlignment"), ("wrap", "wrap")]:
            v = bodyPr.get(attr)
            if v:
                bp[key] = v
        for ins in ["lIns", "rIns", "tIns", "bIns"]:
            v = bodyPr.get(ins)
            if v is not None:
                bp[ins + "_pt"] = emu_to_pt(int(v))
        if bp:
            result["bodyProperties"] = bp

    lstStyle = txBody.find("a:lstStyle", NS)
    if lstStyle is not None:
        levels = {}
        for lvl_el in lstStyle:
            tag = lvl_el.tag.split("}")[-1]
            lvl = {}
            pPr_sub = {}
            for attr, key in [("algn", "alignment"), ("indent", None), ("marL", None)]:
                v = lvl_el.get(attr)
                if v:
                    if key == "alignment":
                        pPr_sub[key] = v
                    elif attr == "indent":
                        pPr_sub["indent_pt"] = emu_to_pt(int(v))
                    elif attr == "marL":
                        pPr_sub["marginLeft_pt"] = emu_to_pt(int(v))
            for tag2, key in [("a:spcBef", "spaceBefore"), ("a:spcAft", "spaceAfter"), ("a:lnSpc", "lineSpacing")]:
                el = lvl_el.find(tag2, NS)
                if el is not None:
                    pts = el.find("a:spcPts", NS)
                    pct = el.find("a:spcPct", NS)
                    if pts is not None:
                        pPr_sub[key + "_pt"] = hundredths_to_pt(int(pts.get("val", 0)))
                    elif pct is not None:
                        pPr_sub[key + "_pct"] = thousandths_to_pct(int(pct.get("val", 0)))
            defRPr = lvl_el.find("a:defRPr", NS)
            if defRPr is not None:
                run = {}
                sz = defRPr.get("sz")
                if sz:
                    run["fontSize_pt"] = hundredths_to_pt(int(sz))
                b = defRPr.get("b")
                if b is not None:
                    run["bold"] = b == "1"
                latin = defRPr.find("a:latin", NS)
                if latin is not None:
                    run["fontFamily"] = latin.get("typeface")
                solidFill = defRPr.find(".//a:solidFill", NS)
                clr = parse_color_element(solidFill)
                if clr:
                    run["color"] = clr
                if run:
                    pPr_sub["defaultRunProps"] = run
            if pPr_sub:
                levels[tag] = pPr_sub
        if levels:
            result["listStyle"] = levels

    para_styles = []
    for para in tf.paragraphs:
        p_style = extract_para_props(para)
        runs_out = []
        for run in para.runs:
            rp = extract_run_props(run)
            if run.text:
                rp["sampleText"] = run.text
            if rp:
                runs_out.append(rp)
        if runs_out:
            p_style["runs"] = runs_out
        if p_style:
            para_styles.append(p_style)
    if para_styles:
        result["paragraphs"] = para_styles

    return result

def extract_shape(shape, location=None):
    data = {
        "shapeName": shape.name,
        "shapeType": str(shape.shape_type),
        "position": {
            "x_pt": emu_to_pt(shape.left),
            "y_pt": emu_to_pt(shape.top),
        },
        "size": {
            "width_pt": emu_to_pt(shape.width),
            "height_pt": emu_to_pt(shape.height),
        },
    }

    try:
        if hasattr(shape, "rotation") and shape.rotation:
            data["rotation"] = shape.rotation
    except Exception:
        pass

    if shape.is_placeholder:
        ph = shape.placeholder_format
        ph_type = str(ph.type).split(".")[-1]
        idx = ph.idx
        if idx == INVALID_PH_IDX:
            inferred_idx = 0 if "TITLE" in ph_type else (
                1 if any(k in ph_type for k in ("BODY", "OBJECT", "TEXT")) else None
            )
            role = "Title" if "TITLE" in ph_type else (
                "Content" if any(k in ph_type for k in ("BODY", "OBJECT", "TEXT")) else "A"
            )
            where = f" on {location}" if location else ""
            msg = (f"{role} placeholder{where} isn't set up correctly in this template "
                   f"(it's missing internal setup info PowerPoint needs). We automatically "
                   f"corrected it so your presentation still converts properly, but for a "
                   f"fully clean template, re-create this box in PowerPoint "
                   f"(View → Slide Master).")
            print(f"  Warning: {msg}")
            _extraction_warnings.append(msg)
            idx = inferred_idx
        data["placeholder"] = {
            "type": ph_type,
            "idx": idx,
        }

    spPr = shape._element.find(".//p:spPr", NS)
    if spPr is None:
        spPr = shape._element.find(".//spPr")
    if spPr is not None:
        solidFill = spPr.find(".//a:solidFill", NS)
        noFill = spPr.find("a:noFill", NS)
        if noFill is not None:
            data["fill"] = "none"
        elif solidFill is not None:
            data["fill"] = parse_color_element(solidFill)
        ln = spPr.find("a:ln", NS)
        if ln is not None:
            border = {}
            w = ln.get("w")
            if w:
                border["width_pt"] = emu_to_pt(int(w))
            lnFill = ln.find(".//a:solidFill", NS)
            lnNoFill = ln.find("a:noFill", NS)
            if lnNoFill is not None:
                border["color"] = "none"
            elif lnFill is not None:
                border["color"] = parse_color_element(lnFill)
            dash = ln.find("a:prstDash", NS)
            if dash is not None:
                border["dashStyle"] = dash.get("val")
            if border:
                data["border"] = border

    if shape.has_text_frame:
        data["textBody"] = extract_text_body(shape.text_frame)

    return {k: v for k, v in data.items() if v is not None}

def extract_theme(prs):
    master_part = prs.slide_master.part
    theme_el = None
    for rel in master_part.rels.values():
        if "theme" in rel.reltype:
            try:
                blob = rel.target_part.blob
                theme_el = etree.fromstring(blob)
                break
            except Exception:
                pass

    if theme_el is None:
        return {}

    colors = {}
    clr_scheme = theme_el.find(".//a:clrScheme", NS)
    if clr_scheme is not None:
        for child in clr_scheme:
            name = child.tag.split("}")[-1]
            clr = parse_color_element(child)
            if clr:
                colors[name] = clr

    fonts = {}
    font_scheme = theme_el.find(".//a:fontScheme", NS)
    if font_scheme is not None:
        major = font_scheme.find("a:majorFont/a:latin", NS)
        minor = font_scheme.find("a:minorFont/a:latin", NS)
        if major is not None:
            fonts["major"] = major.get("typeface")
        if minor is not None:
            fonts["minor"] = minor.get("typeface")

    return {
        "name": theme_el.get("name", ""),
        "colorScheme": colors,
        "fontScheme": fonts,
    }

def extract_master_text_styles(master):
    txStyles = master.element.find("p:txStyles", NS)
    if txStyles is None:
        return {}

    result = {}
    for style_el in txStyles:
        style_name = style_el.tag.split("}")[-1]
        levels = {}
        for lvl_el in style_el:
            tag = lvl_el.tag.split("}")[-1]
            lvl = {}
            algn = lvl_el.get("algn")
            if algn:
                lvl["alignment"] = algn
            indent = lvl_el.get("indent")
            if indent:
                lvl["indent_pt"] = emu_to_pt(int(indent))
            marL = lvl_el.get("marL")
            if marL:
                lvl["marginLeft_pt"] = emu_to_pt(int(marL))
            for tag2, key in [("a:spcBef", "spaceBefore"), ("a:spcAft", "spaceAfter"), ("a:lnSpc", "lineSpacing")]:
                el = lvl_el.find(tag2, NS)
                if el is not None:
                    pts = el.find("a:spcPts", NS)
                    pct = el.find("a:spcPct", NS)
                    if pts is not None:
                        lvl[key + "_pt"] = hundredths_to_pt(int(pts.get("val", 0)))
                    elif pct is not None:
                        lvl[key + "_pct"] = thousandths_to_pct(int(pct.get("val", 0)))
            defRPr = lvl_el.find("a:defRPr", NS)
            if defRPr is not None:
                run = {}
                sz = defRPr.get("sz")
                if sz:
                    run["fontSize_pt"] = hundredths_to_pt(int(sz))
                b = defRPr.get("b")
                if b is not None:
                    run["bold"] = b == "1"
                kern = defRPr.get("kern")
                if kern:
                    run["kerning_pt"] = hundredths_to_pt(int(kern))
                latin = defRPr.find("a:latin", NS)
                if latin is not None:
                    run["fontFamily"] = latin.get("typeface")
                solidFill = defRPr.find(".//a:solidFill", NS)
                clr = parse_color_element(solidFill)
                if clr:
                    run["color"] = clr
                if run:
                    lvl["defaultRunProps"] = run
            buChar_el = lvl_el.find("a:buChar", NS)
            if buChar_el is not None:
                lvl["bulletChar"] = buChar_el.get("char")
            buFont_el = lvl_el.find("a:buFont", NS)
            if buFont_el is not None:
                lvl["bulletFont"] = buFont_el.get("typeface")
            buNone_el = lvl_el.find("a:buNone", NS)
            if buNone_el is not None:
                lvl["bulletNone"] = True
            if lvl:
                levels[tag] = lvl
        if levels:
            result[style_name] = levels

    return result

LAYOUT_TYPE_NAMES = {
    0: "title_slide",
    1: "title_and_content",
    2: "section_header",
    3: "two_content",
    4: "comparison",
    5: "title_only",
    6: "blank",
    7: "content_with_caption",
    8: "picture_with_caption",
}

def extract_layout(layout, idx):
    layout_label = f"the '{LAYOUT_TYPE_NAMES.get(idx, f'layout_{idx}')}' layout"
    placeholders = [extract_shape(s, location=layout_label) for s in layout.placeholders]
    decorative = [extract_shape(s, location=layout_label) for s in layout.shapes if not s.is_placeholder]
    bg = {}
    bg_el = layout.element.find(".//p:bg", NS)
    if bg_el is not None:
        solidFill = bg_el.find(".//a:solidFill", NS)
        if solidFill is not None:
            bg["fill"] = parse_color_element(solidFill)
    result = {
        "layoutName": f"slideLayout{idx + 1}",
        "layoutType": LAYOUT_TYPE_NAMES.get(idx, f"layout_{idx}"),
        "background": bg,
        "placeholders": placeholders,
    }
    if decorative:
        result["decorativeShapes"] = decorative
    return result

def extract_slide(slide, idx):
    shapes = []
    
    try:
        if slide.slide_layout:
            for lay_shape in slide.slide_layout.shapes:
                if not lay_shape.is_placeholder:
                    lay_shape_data = extract_shape(lay_shape)
                    if lay_shape_data:
                        lay_shape_data["isLayoutShape"] = True
                        shapes.append(lay_shape_data)
    except Exception as e:
        print("Failed to extract layout shapes for slide:", idx, e)

    all_slide_shapes = []
    def recurse(container):
        for sh in container:
            if sh.shape_type == 6:
                try:
                    recurse(sh.shapes)
                except Exception:
                    pass
            else:
                all_slide_shapes.append(sh)
    recurse(slide.shapes)
    slide_label = f"template slide {idx + 1}"
    shapes.extend([extract_shape(s, location=slide_label) for s in all_slide_shapes])

    bg_color = None
    try:
        bg_el = slide.element.find(".//p:bg", NS)
        if bg_el is not None:
            solidFill = bg_el.find(".//a:solidFill", NS)
            if solidFill is not None:
                bg_color = parse_color_element(solidFill)
        
        if not bg_color and slide.slide_layout:
            bg_el = slide.slide_layout.element.find(".//p:bg", NS)
            if bg_el is not None:
                solidFill = bg_el.find(".//a:solidFill", NS)
                if solidFill is not None:
                    bg_color = parse_color_element(solidFill)
    except Exception:
        pass

    layout_name = None
    try:
        layout_el = slide.slide_layout.element
        layout_name = layout_el.get("type") or f"slideLayout{slide.slide_layout.slide_master.slide_layouts.index(slide.slide_layout) + 1}"
    except Exception:
        pass

    return {
        "slide_id": str(slide.slide_id) if hasattr(slide, "slide_id") else str(idx + 1),
        "slideNumber": idx + 1,
        "layout": layout_name,
        "backgroundColor": bg_color,
        "shapes": [s for s in shapes if s],
    }

def extract_template(pptx_path):
    _extraction_warnings.clear()
    prs = Presentation(pptx_path)
    master = prs.slide_master

    master_bg = {}
    master_cSld = master.element.find("p:cSld", NS)
    if master_cSld is not None:
        bg_el = master_cSld.find("p:bg", NS)
        if bg_el is not None:
            solidFill = bg_el.find(".//a:solidFill", NS)
            if solidFill is not None:
                master_bg["fill"] = parse_color_element(solidFill)

    style = {
        "_meta": {
            "source": pptx_path,
            "description": "Style guide extracted from PPTX template. Sizes in pt, colors as #hex or scheme:name.",
            "units": "pt for sizes/spacing; #RRGGBB or scheme:<token> for colors",
        },
        "presentation": {
            "width_pt": emu_to_pt(prs.slide_width),
            "height_pt": emu_to_pt(prs.slide_height),
        },
        "theme": extract_theme(prs),
        "masterBackground": master_bg,
        "masterTextStyles": extract_master_text_styles(master),
        "masterShapes": [extract_shape(s) for s in master.shapes],
        "slideLayouts": [extract_layout(lay, i) for i, lay in enumerate(master.slide_layouts)],
        "slides": [extract_slide(sl, i) for i, sl in enumerate(prs.slides)],
    }

    if _extraction_warnings:
        style["_meta"]["warnings"] = list(_extraction_warnings)

    return style
